"""Build the corrected ARISE Kick-off Meeting deck (11 slides).

Visual style closely matches the original deck shared as a PDF:
- Dark navy cover with teal accents
- White content slides with navy titles and a thin teal underline
- Light-fill rounded cards
- Tables with dark navy header rows

Corrected slides are marked with a star (*) before the title so the
reviewer can spot the changes at a glance. Tone and look are preserved.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# --- Palette ---
DARK_BG    = RGBColor(0x0F, 0x1C, 0x2E)
NAVY       = RGBColor(0x12, 0x29, 0x4D)
TEAL       = RGBColor(0x1F, 0x9C, 0x88)
LIGHT_BLUE = RGBColor(0x55, 0xC0, 0xEE)
BLUE_ICON  = RGBColor(0x1B, 0x76, 0xD2)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x8A, 0x8A, 0x8A)
GREY_TXT   = RGBColor(0x55, 0x65, 0x75)
LIGHT_FILL = RGBColor(0xF4, 0xF6, 0xF9)
INK        = RGBColor(0x0F, 0x1C, 0x2E)
RED        = RGBColor(0xD3, 0x2F, 0x2F)
RED_LIGHT  = RGBColor(0xFD, 0xEC, 0xEC)
GREEN      = RGBColor(0x2E, 0x7D, 0x32)
AMBER      = RGBColor(0xE0, 0xA5, 0x16)
BORDER     = RGBColor(0xCF, 0xD8, 0xE2)
WP_GREY    = RGBColor(0x47, 0x5A, 0x6E)
WP_BLUE    = RGBColor(0x3A, 0x8D, 0xD3)
WP_PURPLE  = RGBColor(0x84, 0x52, 0xC6)
WP_ORANGE  = RGBColor(0xE6, 0x95, 0x2A)
WP_GREEN   = RGBColor(0x4A, 0xAE, 0x6E)
WP_PINK    = RGBColor(0xC2, 0x6B, 0x9A)  # for Bridge

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

STAR = "★"


def prs_setup():
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def text(slide, x, y, w, h, content, *, size=12, bold=False, italic=False,
         color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    lines = content if isinstance(content, list) else [content]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def rect(slide, x, y, w, h, *, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, int(x), int(y), int(w), int(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def rounded_card(slide, x, y, w, h, *, fill=LIGHT_FILL, line=BORDER, line_w=0.5,
                 top_band=None, top_band_h=None):
    """Light-fill rounded card. Optional top accent band."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h))
    card.adjustments[0] = 0.06
    card.fill.solid(); card.fill.fore_color.rgb = fill
    card.line.color.rgb = line
    card.line.width = Pt(line_w)
    card.shadow.inherit = False
    if top_band:
        band_h = top_band_h or Inches(0.06)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      int(x), int(y), int(w), int(band_h))
        band.fill.solid(); band.fill.fore_color.rgb = top_band
        band.line.fill.background()
        band.shadow.inherit = False
    return card


def slide_title(slide, title, *, starred=False):
    """Standard content-slide title block: large navy title + thin teal rule."""
    if starred:
        title = f"{STAR} {title}"
    text(slide, Inches(0.55), Inches(0.35), Inches(12.3), Inches(0.7),
         title, size=30, bold=True, color=NAVY)
    rect(slide, Inches(0.55), Inches(1.05), Inches(12.2), Emu(20000),
         fill=TEAL)


def table_simple(slide, x, y, w, col_widths, header, rows,
                 *, header_fill=NAVY, header_color=WHITE,
                 row_alt=LIGHT_FILL, font_size=10, header_size=10,
                 first_col_bold=False, line_color=BORDER):
    """Render a simple table with a dark header row and alternating row fills."""
    assert sum(col_widths) <= 1.001
    n_cols = len(header)
    assert len(col_widths) == n_cols
    n_rows = len(rows) + 1  # +1 header
    row_h = Inches(0.42)
    header_h = Inches(0.45)
    yy = y
    # Header
    xx = x
    for i, ttl in enumerate(header):
        cw = int(w * col_widths[i])
        cell = rect(slide, xx, yy, cw, header_h, fill=header_fill)
        tb = text(slide, xx + Inches(0.1), yy, cw - Inches(0.2), header_h,
                  ttl, size=header_size, bold=True, color=header_color,
                  anchor=MSO_ANCHOR.MIDDLE)
        xx += cw
    yy += header_h
    # Body
    for ri, row in enumerate(rows):
        xx = x
        rfill = LIGHT_FILL if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cw = int(w * col_widths[ci])
            cell = rect(slide, xx, yy, cw, row_h, fill=rfill, line=line_color, line_w=0.25)
            tb = text(slide, xx + Inches(0.1), yy, cw - Inches(0.2), row_h,
                      str(val), size=font_size,
                      bold=(first_col_bold and ci == 0),
                      color=INK, anchor=MSO_ANCHOR.MIDDLE)
            xx += cw
        yy += row_h


def code_chip(slide, x, y, w, h, text_val, *, fill=BLUE_ICON, color=WHITE,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    chip = slide.shapes.add_shape(shape, int(x), int(y), int(w), int(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        chip.adjustments[0] = 0.3
    chip.fill.solid(); chip.fill.fore_color.rgb = fill
    chip.line.fill.background()
    chip.shadow.inherit = False
    text(slide, x, y, w, h, text_val, size=10, bold=True,
         color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return chip


# === Slide builders ===

def slide1_cover(prs):
    s = blank(prs)
    set_bg(s, DARK_BG)

    # Teal highlight bar behind ARISE
    bar = slide_bar(s, Inches(5.2), Inches(2.55), Inches(3.0), Inches(1.4), TEAL)
    bar.fill.transparency = 0.7  # not supported in python-pptx easily, leave full

    text(s, Inches(0.5), Inches(2.1), Inches(12.3), Inches(0.4),
         "KICK-OFF MEETING", size=16, bold=True, color=LIGHT_BLUE,
         align=PP_ALIGN.CENTER, font="Calibri")
    text(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5),
         "ARISE", size=88, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.5),
         "Augmented Rehabilitation & Intelligent System for Enhancement",
         size=22, color=GRAY, align=PP_ALIGN.CENTER, font="Calibri")
    text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.4),
         "The first proactive \"AI Coach\" for Sit-to-Stand (STS) rehabilitation",
         size=16, color=WHITE, align=PP_ALIGN.CENTER, italic=False)

    # Date box
    bx = Inches(4.3); by = Inches(5.55); bw = Inches(4.7); bh = Inches(1.5)
    box = slide_bar(s, bx, by, bw, bh, DARK_BG, line=TEAL, line_w=2)
    text(s, bx, by + Inches(0.15), bw, Inches(0.35),
         "Start Date:", size=14, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, bx, by + Inches(0.5), bw, Inches(0.4),
         "June 18, 2026", size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    text(s, bx, by + Inches(0.92), bw, Inches(0.3),
         "Expected End:", size=14, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, bx, by + Inches(1.18), bw, Inches(0.35),
         "December 18, 2027", size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


def slide_bar(slide, x, y, w, h, fill, *, line=None, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               int(x), int(y), int(w), int(h))
    s.adjustments[0] = 0.12
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def slide2_problem_vision(prs):
    s = blank(prs); set_bg(s, WHITE)
    slide_title(s, "The Problem and the Vision")

    # Left column
    xL = Inches(0.55); xR = Inches(7.5)
    # Problem
    icon_box(s, xL, Inches(1.5), Inches(0.45), Inches(0.45), RED, "!")
    text(s, xL + Inches(0.6), Inches(1.5), Inches(6.6), Inches(0.4),
         "The Problem:", size=14, bold=True, color=INK)
    text(s, xL + Inches(0.6), Inches(1.95), Inches(6.6), Inches(1.0),
         "Discontinuity of care in motor rehabilitation (STS). "
         "Incorrect implementation at home and lack of objective data.",
         size=12, color=GREY_TXT)
    # Solution
    icon_box(s, xL, Inches(3.1), Inches(0.45), Inches(0.45), TEAL, "+")
    text(s, xL + Inches(0.6), Inches(3.1), Inches(6.6), Inches(0.4),
         "The Solution:", size=14, bold=True, color=INK)
    text(s, xL + Inches(0.6), Inches(3.55), Inches(6.6), Inches(1.2),
         "A product/service that acts as a \"Virtual Therapist.\" "
         "Not just reactive estimation, but 3D measurement and proactive guidance (forecasting).",
         size=12, color=GREY_TXT)
    # Pillars
    icon_box(s, xL, Inches(5.0), Inches(0.45), Inches(0.45), TEAL, "#")
    text(s, xL + Inches(0.6), Inches(5.0), Inches(6.6), Inches(0.4),
         "The 3 Pillars of ARISE:", size=14, bold=True, color=INK)
    text(s, xL + Inches(0.7), Inches(5.45), Inches(6.5), Inches(1.7),
         ["1.  Core AI (Measurement and Prediction)",
          "2.  \"Coach\" Edge Device",
          "3.  Clinical Cloud Dashboard"],
         size=12, color=GREY_TXT)

    # Right image placeholder
    rounded_card(s, xR, Inches(1.5), Inches(5.3), Inches(5.0),
                 fill=LIGHT_FILL, line=BORDER)
    text(s, xR, Inches(3.7), Inches(5.3), Inches(0.5),
         "[Photo of patient + therapist]",
         size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


def icon_box(slide, x, y, w, h, fill, glyph):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               int(x), int(y), int(w), int(h))
    s.adjustments[0] = 0.25
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    text(slide, x, y, w, h, glyph, size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide3_objectives(prs):
    s = blank(prs); set_bg(s, WHITE)
    slide_title(s, "Project Objectives (TRL6)")

    cards = [
        ("Patients",
         "Improve ADL and reduce fall risk by providing intuitive visual "
         "and auditory biofeedback, designed with low cognitive load."),
        ("Therapist",
         "Provide an advanced decision dashboard based on objective kinematic "
         "data (not qualitative estimates) to personalise treatment plans."),
        ("Business (EaaS)",
         "Validate the technology to launch the \"Equipment as a Service\" (EaaS) "
         "business model for clinics, nursing homes, and home tele-rehabilitation."),
    ]
    card_w = Inches(3.95); gap = Inches(0.25)
    x0 = Inches(0.55) + (Inches(12.2) - 3 * card_w - 2 * gap) / 2
    y = Inches(2.4); h = Inches(4.4)
    for i, (title, body) in enumerate(cards):
        x = x0 + i * (card_w + gap)
        rounded_card(s, x, y, card_w, h, fill=LIGHT_FILL, line=BORDER)
        # Icon stub
        icon_box(s, x + (card_w - Inches(0.7)) / 2, y + Inches(0.5),
                 Inches(0.7), Inches(0.7), BLUE_ICON, "*")
        text(s, x, y + Inches(1.5), card_w, Inches(0.5),
             title, size=20, bold=True, color=INK,
             align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.3), y + Inches(2.15), card_w - Inches(0.6),
             Inches(2.0), body,
             size=12, color=GREY_TXT, align=PP_ALIGN.CENTER)


def slide4_gantt(prs):
    """Gantt chart with the new Bridge phase (M8-M18) added."""
    s = blank(prs); set_bg(s, WHITE)
    slide_title(s, "Work and Development Plan (GANTT)", starred=True)
    text(s, Inches(0.55), Inches(1.2), Inches(12.2), Inches(0.4),
         "18-Month Roadmap: 18 June 2026 - 18 December 2027",
         size=13, color=TEAL, italic=True)

    # Outer chart frame
    cx = Inches(0.55); cy = Inches(1.8); cw = Inches(12.2); ch = Inches(5.0)
    rounded_card(s, cx, cy, cw, ch, fill=LIGHT_FILL, line=BORDER)

    # Quarter labels (T1..T6)
    qlabel_y = cy + Inches(0.15)
    text(s, cx + Inches(2.5), qlabel_y, Inches(1.5), Inches(0.3),
         "T1 ('26)", size=10, bold=True, color=INK)
    text(s, cx + Inches(4.0), qlabel_y, Inches(1.5), Inches(0.3),
         "T2 ('26)", size=10, bold=True, color=INK)
    text(s, cx + Inches(5.5), qlabel_y, Inches(1.5), Inches(0.3),
         "T3 ('26/'27)", size=10, bold=True, color=INK)
    text(s, cx + Inches(7.6), qlabel_y, Inches(1.5), Inches(0.3),
         "T4 ('27)", size=10, bold=True, color=INK)
    text(s, cx + Inches(9.1), qlabel_y, Inches(1.5), Inches(0.3),
         "T5 ('27)", size=10, bold=True, color=INK)
    text(s, cx + Inches(10.6), qlabel_y, Inches(1.5), Inches(0.3),
         "T6 ('27)", size=10, bold=True, color=INK)

    # Month columns (7..12, 1..12)
    # Use the existing 18 months: Jul'26 to Dec'27 (7,8,9,10,11,12,1..12)
    months = [7, 8, 9, 10, 11, 12, 1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12]
    n = len(months)
    col_w = (cw - Inches(2.0)) / n
    col_x0 = cx + Inches(2.0)
    label_y = cy + Inches(0.5)
    for i, m in enumerate(months):
        x = col_x0 + i * col_w
        text(s, x, label_y, col_w, Inches(0.3),
             str(m), size=8, color=GREY_TXT, align=PP_ALIGN.CENTER)
        # vertical guide
        guide = rect(s, x + col_w / 2, cy + Inches(0.85),
                     Emu(8000), Inches(4.0), fill=BORDER)

    # Header dividing rule
    rect(s, cx, cy + Inches(0.85), cw, Emu(15000), fill=BORDER)

    # WP rows
    wp_rows = [
        ("WP1: Project Mngt.",        0, 17, "M1-M18 (Governance & Ethics)", WP_GREY),
        ("WP2: Analysis & Architecture.", 1, 6,  "M2-M7 (Cloud & UX/UI)",    WP_BLUE),
        ("WP3: AI & HW Research",     3, 10, "M4-M11 (Dataset & Pose 3D)",   WP_PURPLE),
        ("Bridge (AI continuation)",  7, 17, "M8-M18 (AI refinement during integration & validation)", WP_PINK),
        ("WP4: Platform Dev",         8, 13, "M9-M14 (Edge & Dashboard)",    WP_ORANGE),
        ("WP5: TRL6 Validation",     13, 17, "M14-M18 (Real Environment Test)", WP_GREEN),
    ]
    row_y0 = cy + Inches(1.0); row_h = Inches(0.55)
    for ri, (label, start_i, end_i, btext, color) in enumerate(wp_rows):
        ry = row_y0 + ri * row_h
        text(s, cx + Inches(0.1), ry, Inches(1.9), row_h,
             label, size=9, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        # bar
        bx = col_x0 + start_i * col_w
        bw = (end_i - start_i + 1) * col_w
        bar_h = Inches(0.32) if ri != 3 else Inches(0.18)  # Bridge thinner
        bar_y = ry + (row_h - bar_h) / 2
        bar = slide_bar(s, bx, bar_y, bw, bar_h, color)
        text(s, bx + Inches(0.05), bar_y, bw - Inches(0.1), bar_h,
             btext, size=8, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Milestone diamonds (D1.1 at M3, D1.2 at M4, D2.1 at M5, D3.1 at M11,
    # D4.1 at M14, M5.1 at M15, T5.4 at M16, D5.1 at M18, D1.3 at M18)
    def diamond(idx, label, color=AMBER, y=cy + Inches(1.0) + Inches(0.55) * 0):
        x = col_x0 + idx * col_w + col_w / 2 - Inches(0.14)
        d = slide.shapes.add_shape if False else None
    # We'll just place small chips above the bars instead, simpler
    def chip(month_idx, label, y_row_i):
        x = col_x0 + month_idx * col_w + col_w / 2 - Inches(0.4)
        y = row_y0 + y_row_i * row_h - Inches(0.18)
        chip_shape = slide_bar(s, x, y, Inches(0.8), Inches(0.25), AMBER)
        text(s, x, y, Inches(0.8), Inches(0.25),
             label, size=8, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    chip(1, "D1.1", 0)   # ~M3 (Aug)
    chip(2, "D1.2", 0)   # ~M4 (Sep)
    chip(4, "D2.1", 1)   # ~M5 (Nov)
    chip(9, "D3.1", 2)   # ~M11 (Apr)
    chip(13, "D4.1", 4)  # ~M14 (Aug)
    chip(14, "M5.1", 5)  # ~M15 (Sep)
    chip(15, "T5.4", 5)  # ~M16 (Oct)
    chip(17, "D5.1", 5)  # ~M18 (Dec)
    chip(17, "D1.3", 0)  # ~M18 (Dec)


def slide5_deliverables(prs):
    """Deliverables list with the missing sub-deliverables added."""
    s = blank(prs); set_bg(s, WHITE)
    slide_title(s, "Project Deliverables and Milestones List", starred=True)

    header = ["Code", "Title", "Operational Description"]
    rows = [
        ("D1.1", "Clinical Requirements and Biomechanical KPIs",
         "Defines the scientific Ground Truth and the KPIs (angles, velocities) needed for AI training."),
        ("D1.2", "Submission to the Ethics Committee",
         "Drafting of the clinical protocol, inclusion criteria, and formal submission for CER approval."),
        ("D1.3", "Final Report and Dissemination",
         "Project management report and reporting of scientific dissemination activities carried out."),
        ("D2.1", "UI/UX Design (Mockup)",
         "Creation and validation of visual prototypes for the patient (Biofeedback) and therapist (Dashboard) interfaces."),
        ("T3.2", "STS Dataset Acquisition and Annotation",
         "Acquisition of STS recordings at the clinical site and DINOGMI-led double-blind annotation, building the AI training corpus."),
        ("D3.1", "AI Models and STS Datasets",
         "Clinical data acquisition and training of algorithms for real-time 3D estimation and forecasting."),
        ("D4.1", "TRL6 Platform Prototype",
         "Release of the Coach hardware device and cloud infrastructure, ready for the real world."),
        ("D4.2", "Cloud Platform Release",
         "Cloud back-end, ingestion API, and Clinical Dashboard delivered as a separate release distinct from the Coach edge hardware."),
        ("M5.1", "Start of Trial (TRL6)",
         "Clinical setup, patient recruitment, and official start of testing inside the gym."),
        ("T5.4", "AI vs. Gold Standard Comparison",
         "Statistical analysis of AI metrics compared with laboratory optoelectronic systems (REHELAB)."),
        ("D5.1", "TRL6 Validation Report",
         "Final document on scientific validation, usability, patient adherence, and clinical impact."),
        ("D5.2", "Patient and Therapist Feedback Report",
         "Intermediate usability and adherence report at the mid-point of WP5, feeding D5.1 with early signals."),
    ]

    y0 = Inches(1.4)
    table_simple(s, Inches(0.55), y0, Inches(12.2),
                 col_widths=[0.10, 0.32, 0.58],
                 header=header, rows=[(r[1], r[2]) for r in rows],
                 font_size=9, header_size=11)
    # Re-draw the code chips on column 1
    chip_w = Inches(0.85); chip_h = Inches(0.32)
    yy = y0 + Inches(0.45)
    row_h = Inches(0.42)
    for code, _, _ in rows:
        is_milestone = code.startswith("M")
        is_task = code.startswith("T")
        fill = (GREEN if is_milestone else (WP_GREY if is_task else BLUE_ICON))
        cx = Inches(0.55) + Inches(0.1)
        chip_y = yy + (row_h - chip_h) / 2
        code_chip(s, cx, chip_y, chip_w, chip_h, code, fill=fill)
        yy += row_h


def slide6_unige(prs):
    s = blank(prs); set_bg(s, WHITE)
    slide_title(s, "Focus Partner: University of Genoa (DINOGMI)", starred=True)

    # Left card: role + critical + mitigation
    lx = Inches(0.55); ly = Inches(1.4); lw = Inches(4.5); lh = Inches(5.6)
    rounded_card(s, lx, ly, lw, lh, fill=WHITE, line=BORDER)
    y = ly + Inches(0.25)
    text(s, lx + Inches(0.3), y, lw - Inches(0.6), Inches(0.4),
         "Role & Activities", size=14, bold=True, color=INK)
    rect(s, lx + Inches(0.3), y + Inches(0.42), lw - Inches(0.6), Emu(10000),
         fill=BORDER)
    text(s, lx + Inches(0.3), y + Inches(0.6), lw - Inches(0.6), Inches(0.3),
         "Scientific Partner:", size=11, bold=True, color=INK)
    text(s, lx + Inches(0.3), y + Inches(0.9), lw - Inches(0.6), Inches(0.8),
         "Member of the Biomechanical Validation Committee and definition of the Gold Standard at REHELAB.",
         size=11, color=GREY_TXT)
    text(s, lx + Inches(0.3), y + Inches(1.85), lw - Inches(0.6), Inches(0.3),
         "Initial Phase:", size=11, bold=True, color=INK)
    text(s, lx + Inches(0.3), y + Inches(2.15), lw - Inches(0.6), Inches(0.9),
         "Define the scientific Ground Truth (biomechanical KPIs) by mid-August 2026, "
         "and only then start drafting the Ethics Committee protocol, since the KPIs are what the trial measures.",
         size=11, color=GREY_TXT)
    # Critical issues + mitigation (red box)
    crit_y = y + Inches(3.3); crit_h = Inches(2.05)
    rect(s, lx + Inches(0.3), crit_y, Inches(0.05), crit_h, fill=RED)
    rounded_card(s, lx + Inches(0.35), crit_y, lw - Inches(0.7), crit_h,
                 fill=RED_LIGHT, line=RED_LIGHT)
    text(s, lx + Inches(0.5), crit_y + Inches(0.1), lw - Inches(1.0), Inches(0.3),
         "Critical issues:", size=11, bold=True, color=RED)
    text(s, lx + Inches(0.5), crit_y + Inches(0.4), lw - Inches(1.0), Inches(0.6),
         "Slow approval by the Regional Ethics Committee, and KPI sign-off slipping past M3.",
         size=10, color=INK)
    text(s, lx + Inches(0.5), crit_y + Inches(1.1), lw - Inches(1.0), Inches(0.3),
         "Mitigation:", size=11, bold=True, color=RED)
    text(s, lx + Inches(0.5), crit_y + Inches(1.4), lw - Inches(1.0), Inches(0.6),
         "Lock D1.1 KPIs by mid-August, then run protocol drafting and EC pre-submission "
         "in parallel from M3 to M4 to absorb committee review time.",
         size=10, color=INK)

    # Right: deliverables table
    rx = Inches(5.3); ry = Inches(1.4); rw = Inches(7.5)
    table_simple(s, rx, ry, rw,
                 col_widths=[0.15, 0.55, 0.30],
                 header=["ID", "Deliverable / Activity", "Expected Deadline"],
                 rows=[("D1.1", "Definition of Clinical Requirements and KPIs", "Mid-August 2026"),
                       ("D1.2", "Submission to the Ethics Committee",          "Mid-September 2026"),
                       ("T5.4", "AI vs. Gold Standard Comparison",              "Mid-October 2027"),
                       ("D5.1", "Final Report - Validation & Dissemination",   "18 December 2027")],
                 font_size=10, header_size=11)


def slide7_buccarella(prs):
    s = blank(prs); set_bg(s, WHITE)
    slide_title(s, "Focus Partner: Buccarella Associates")

    lx = Inches(0.55); ly = Inches(1.4); lw = Inches(4.5); lh = Inches(5.6)
    rounded_card(s, lx, ly, lw, lh, fill=WHITE, line=BORDER)
    y = ly + Inches(0.25)
    text(s, lx + Inches(0.3), y, lw - Inches(0.6), Inches(0.4),
         "Role & Activities", size=14, bold=True, color=INK)
    rect(s, lx + Inches(0.3), y + Inches(0.42), lw - Inches(0.6), Emu(10000),
         fill=BORDER)
    text(s, lx + Inches(0.3), y + Inches(0.6), lw - Inches(0.6), Inches(0.3),
         "Clinical-Operational Partner:", size=11, bold=True, color=INK)
    text(s, lx + Inches(0.3), y + Inches(0.9), lw - Inches(0.6), Inches(0.5),
         "Deputy for Real-World Validation.", size=11, color=GREY_TXT)
    text(s, lx + Inches(0.3), y + Inches(1.4), lw - Inches(0.6), Inches(1.6),
         "Unlike the University, which deals with pure scientific validation, "
         "the Buccarella Studio provides the fundamental perspective of the real environment.",
         size=11, bold=True, color=INK)

    crit_y = y + Inches(3.3); crit_h = Inches(2.05)
    rect(s, lx + Inches(0.3), crit_y, Inches(0.05), crit_h, fill=RED)
    rounded_card(s, lx + Inches(0.35), crit_y, lw - Inches(0.7), crit_h,
                 fill=RED_LIGHT, line=RED_LIGHT)
    text(s, lx + Inches(0.5), crit_y + Inches(0.1), lw - Inches(1.0), Inches(0.3),
         "Critical issues:", size=11, bold=True, color=RED)
    text(s, lx + Inches(0.5), crit_y + Inches(0.4), lw - Inches(1.0), Inches(0.6),
         "Difficulty in recruiting patients during the months of WP5.",
         size=10, color=INK)
    text(s, lx + Inches(0.5), crit_y + Inches(1.1), lw - Inches(1.0), Inches(0.3),
         "Mitigation:", size=11, bold=True, color=RED)
    text(s, lx + Inches(0.5), crit_y + Inches(1.4), lw - Inches(1.0), Inches(0.6),
         "Pre-screening and pre-enrollment in the months preceding the tests.",
         size=10, color=INK)

    rx = Inches(5.3); ry = Inches(1.4); rw = Inches(7.5)
    table_simple(s, rx, ry, rw,
                 col_widths=[0.15, 0.55, 0.30],
                 header=["ID", "Deliverable / Activity", "Expected Deadline"],
                 rows=[("D1.2", "Workflow Protocol for the Ethics Committee", "Mid-September 2026"),
                       ("D2.1", "Feedback UI/UX - Biofeedback and Dashboard", "Mid-November 2026"),
                       ("T3.2", "STS Dataset Acquisition (clinical hosting)", "M9 to M15"),
                       ("M5.1", "Experimental Launch - TRL6 Test",            "Mid-September 2027"),
                       ("D5.1", "Validation Report - Data and Adherence",     "18 December 2027")],
                 font_size=10, header_size=11)


def slide8_integrated(prs):
    s = blank(prs); set_bg(s, WHITE)
    slide_title(s, "Integrated Operational Contribution (UniGe & Buccarella)", starred=True)

    table_simple(s, Inches(0.55), Inches(1.4), Inches(12.2),
                 col_widths=[0.08, 0.20, 0.13, 0.30, 0.29],
                 header=["ID", "Deliverable", "Expiration",
                         "UniGe Contribution (Scientific)",
                         "Studio Buccarella Contribution (Clinical)"],
                 rows=[
                     ("D1.1", "Clinical Requirements and KPIs", "Mid-Aug '26",
                      "Defines the scientific Ground Truth and KPIs (angles, velocity) used to train the AI.",
                      "Supplies the functional and usability requirements of the patients."),
                     ("D1.2", "Ethics Protocol", "Mid-Sep '26",
                      "Scientific and formal coordination for drafting the CER protocol, after D1.1 KPI sign-off.",
                      "Practical definition of patient inclusion criteria, logistics, and clinical workflow."),
                     ("D2.1", "Feedback UI/UX", "Mid-Nov '26",
                      "Consulted on the clinical and scientific accuracy of the kinematic data shown in the mockups.",
                      "Tests intuitiveness of biofeedback for elderly patients and usefulness of the dashboard."),
                     ("T3.2", "STS Dataset Acquisition", "M9 to M15",
                      "Methodological supervision of the acquisition and DINOGMI-led ground-truth annotation.",
                      "Provides the clinical site, patients, and healthy volunteers for the recordings."),
                     ("M5.1", "Setup Startup (TRL6)", "Mid-Sep '27",
                      "Methodological supervision of the setup to ensure the validity of the data in the environment.",
                      "Patient recruitment, consent management, and session management."),
                     ("T5.4", "AI vs. Gold Std", "Mid-Oct '27",
                      "Statistical comparison of ARISE measurements vs the REHELAB optoelectronic Gold Standard.",
                      "Provides field clinical logs to support and explain any detected anomalies."),
                     ("D5.1", "TRL6 Final Report", "18 Dec '27",
                      "Final scientific validation, statistical analysis, academic paper writing.",
                      "Data on patient adherence, practical usability, impact on workload."),
                 ],
                 font_size=9, header_size=10)


def slide9_raci(prs):
    """RACI matrix with corrections: D2.1 UniGe I->C, M5.1 UniGe C->R, plus T3.2 row."""
    s = blank(prs); set_bg(s, WHITE)
    slide_title(s, "RACI Matrix of Clinical Activities and Validation", starred=True)

    # Render the table without the RACI columns first to add chips on top
    rows = [
        ("D1.1", "Clinical Requirements and Biomechanical KPIs", "A", "R", "C"),
        ("D1.2", "Protocol and Submission to the Ethics Committee", "A", "R", "R"),
        ("D2.1", "Feedback UI/UX (Mockup Coach & Dashboard)", "A/R", "C", "C"),
        ("T3.2", "STS Dataset Acquisition and Annotation", "A", "C", "R"),
        ("M5.1", "Patient Trial Setup and Startup (TRL6)", "A", "R", "R"),
        ("T5.4", "AI Data Comparison vs. Gold Standard", "A", "R", "I"),
        ("D5.1", "TRL6 Final Validation Report", "A/R", "R", "R"),
    ]
    table_simple(s, Inches(0.55), Inches(1.4), Inches(12.2),
                 col_widths=[0.09, 0.43, 0.16, 0.16, 0.16],
                 header=["ID", "Deliverable / Activity",
                         "Innovina S.r.l.", "UniGe (DINOGMI)", "Studio Buccarella"],
                 rows=[(r[1], "", "", "") for r in rows],
                 font_size=10, header_size=11)
    # Overlay chips and code chips
    yy = Inches(1.4) + Inches(0.45)
    row_h = Inches(0.42)
    chip_colors = {"R": GREEN, "A": RED, "C": AMBER, "I": GREY_TXT, "A/R": RED}
    for ri, (code, desc, ic, uc, bc) in enumerate(rows):
        # code chip
        x_code = Inches(0.55) + Inches(0.1)
        code_chip(s, x_code, yy + (row_h - Inches(0.32)) / 2,
                  Inches(0.7), Inches(0.32),
                  code, fill=BLUE_ICON if not code.startswith("T") and not code.startswith("M")
                          else (GREEN if code.startswith("M") else WP_GREY))
        # RACI cells
        col_w = Inches(12.2)
        # col positions for last three columns (widths 0.16, 0.16, 0.16; cum from 0.52 onwards)
        col_x = [
            Inches(0.55) + col_w * 0.52,
            Inches(0.55) + col_w * 0.68,
            Inches(0.55) + col_w * 0.84,
        ]
        col_cell_w = col_w * 0.16
        for ci, code_val in enumerate([ic, uc, bc]):
            cx = col_x[ci]
            chip_w = Inches(0.55); chip_h = Inches(0.32)
            cw_x = cx + (col_cell_w - chip_w) / 2
            cw_y = yy + (row_h - chip_h) / 2
            code_chip(s, cw_x, cw_y, chip_w, chip_h,
                      code_val, fill=chip_colors.get(code_val, GREY_TXT))
        yy += row_h

    # Legend below
    leg_y = Inches(6.0)
    legend_items = [
        ("R", "Responsible — physically carries out the work", GREEN),
        ("A", "Accountable — approves and has final responsibility", RED),
        ("C", "Consulted — provides essential input", AMBER),
        ("I", "Informed — kept informed of progress", GREY_TXT),
    ]
    rounded_card(s, Inches(0.55), leg_y, Inches(12.2), Inches(0.8),
                 fill=LIGHT_FILL, line=BORDER)
    col_w = Inches(12.2) / 4
    for i, (letter, desc, color) in enumerate(legend_items):
        x = Inches(0.55) + i * col_w + Inches(0.2)
        code_chip(s, x, leg_y + Inches(0.18), Inches(0.4), Inches(0.4),
                  letter, fill=color)
        text(s, x + Inches(0.55), leg_y + Inches(0.13), col_w - Inches(0.8),
             Inches(0.6), desc, size=9, color=INK, anchor=MSO_ANCHOR.MIDDLE)


def slide10_strategic(prs):
    s = blank(prs); set_bg(s, WHITE)
    slide_title(s, "Strategic Alignment & Risk Management", starred=True)

    # Critical milestone banner — updated to reflect feedback #3
    bx = Inches(0.55); by = Inches(1.4); bw = Inches(12.2); bh = Inches(1.2)
    rect(s, bx, by, Inches(0.12), bh, fill=LIGHT_BLUE)
    rounded_card(s, bx + Inches(0.12), by, bw - Inches(0.12), bh,
                 fill=DARK_BG, line=DARK_BG)
    text(s, bx + Inches(0.3), by + Inches(0.13), bw - Inches(0.5), Inches(0.4),
         "CRITICAL MILESTONE: TRL6 PROTOTYPE DEPLOYMENT (M14 to M18)",
         size=14, bold=True, color=LIGHT_BLUE)
    text(s, bx + Inches(0.3), by + Inches(0.55), bw - Inches(0.5), Inches(0.6),
         ["Platform availability: Mid-August 2027 (M14)  |  First patient: Mid-September 2027 (M15)  |  AI vs Gold Std: Mid-October 2027 (M16)  |  Final report: 18 December 2027 (M18)",
          "Compressed window: 4 months from platform availability to final report. Risk of recruitment / analysis spillover, mitigated by pre-screening from M11 and a 2-month D5.2 interim report (mid-November 2027)."],
         size=10, color=WHITE)

    # Guide questions left + right
    lq_x = Inches(0.55); rq_x = Inches(6.85); q_y = Inches(2.85)
    q_w = Inches(6.05); q_h = Inches(4.1)

    rounded_card(s, lq_x, q_y, q_w, q_h, fill=LIGHT_FILL, line=BORDER,
                 top_band=LIGHT_BLUE, top_band_h=Inches(0.06))
    text(s, lq_x + Inches(0.3), q_y + Inches(0.2), q_w - Inches(0.6), Inches(0.4),
         "Guide Questions for DINOGMI", size=14, bold=True, color=BLUE_ICON)
    questions_L = [
        ("Ethical Bureaucracy:",
         "What are the worst-case timelines for the CER? "
         "Given the protocol drafting starts after KPI sign-off at M3, what is the plan B if approval slips past M6?"),
        ("AI Feasibility:",
         "How do we ensure that the desk-defined biomechanical KPIs (D1.1) "
         "can be robustly extracted from markerless sensors?"),
        ("Bridge Phase:",
         "During the Bridge phase (M8 to M18), how do we coordinate AI iteration "
         "with WP4 integration and WP5 validation so that issues found in the trial can be triaged without invalidating it?"),
    ]
    qy = q_y + Inches(0.75)
    for label, body in questions_L:
        text(s, lq_x + Inches(0.3), qy, q_w - Inches(0.6), Inches(0.3),
             label, size=11, bold=True, color=INK)
        text(s, lq_x + Inches(0.3), qy + Inches(0.3), q_w - Inches(0.6), Inches(0.85),
             body, size=10, color=GREY_TXT)
        qy += Inches(1.15)

    rounded_card(s, rq_x, q_y, q_w, q_h, fill=LIGHT_FILL, line=BORDER,
                 top_band=TEAL, top_band_h=Inches(0.06))
    text(s, rq_x + Inches(0.3), q_y + Inches(0.2), q_w - Inches(0.6), Inches(0.4),
         "Guide Questions for Buccarella Studio", size=14, bold=True, color=TEAL)
    questions_R = [
        ("Patient Recruitment:",
         "Given a 4-month WP5 window (M14 to M18), what is a realistic estimate of "
         "eligible patients enrollable, and when should pre-screening start (M11?)"),
        ("Workflow Integration:",
         "What is the maximum acceptable time (minutes) to start the Coach in the gym, "
         "so as not to hinder the therapists?"),
        ("Geriatric Usability:",
         "How do we calibrate biofeedback (visual/auditory) to stimulate self-correction "
         "without frightening or confusing an elderly patient?"),
    ]
    qy = q_y + Inches(0.75)
    for label, body in questions_R:
        text(s, rq_x + Inches(0.3), qy, q_w - Inches(0.6), Inches(0.3),
             label, size=11, bold=True, color=INK)
        text(s, rq_x + Inches(0.3), qy + Inches(0.3), q_w - Inches(0.6), Inches(0.85),
             body, size=10, color=GREY_TXT)
        qy += Inches(1.15)


def slide11_next_steps(prs):
    s = blank(prs); set_bg(s, WHITE)
    slide_title(s, "Next Steps (First 30-60 Days)", starred=True)

    boxes = [
        ("1", "Formal Kick-Off Setup", LIGHT_BLUE,
         "Project management tools and document-sharing setup by Innovina. "
         "18-month GANTT alignment (Start: 18 June 2026)."),
        ("2", "Technical Tables - KPI Definition", LIGHT_BLUE,
         "Immediate and priority launch of technical tables (Innovina & DINOGMI) "
         "for the drafting of D1.1 (Clinical Requirements and Biomechanical KPIs)."),
        ("3", "Ethics Committee Protocol (sequenced)", RED,
         "Joint Innovina-DINOGMI-Buccarella drafting of the Clinical Protocol "
         "starts only AFTER D1.1 KPI sign-off (mid-August 2026), to avoid drafting in the dark. "
         "EC pre-submission window M3 to M4."),
        ("4", "Design Startup (WP2)", LIGHT_BLUE,
         "Launch of cloud architecture and UI/UX analysis, gathering initial clinical "
         "feedback on user interaction and dashboard needs."),
    ]
    card_w = Inches(5.95); card_h = Inches(2.5); gap_x = Inches(0.3); gap_y = Inches(0.35)
    x0 = Inches(0.55); y0 = Inches(1.4)
    for i, (num, title, accent, body) in enumerate(boxes):
        col = i % 2; row = i // 2
        x = x0 + col * (card_w + gap_x)
        y = y0 + row * (card_h + gap_y)
        rounded_card(s, x, y, card_w, card_h, fill=LIGHT_FILL, line=BORDER,
                     top_band=accent, top_band_h=Inches(0.10))
        # Number chip
        is_priority = (accent == RED)
        chip_fill = RED if is_priority else BLUE_ICON
        chip_glyph = "!" if is_priority else num
        # Always use a numbered chip — fixes the original 1,2,!,4 bug
        code_chip(s, x + Inches(0.3), y + Inches(0.4),
                  Inches(0.4), Inches(0.4), num, fill=chip_fill,
                  shape=MSO_SHAPE.OVAL)
        # Title + body
        title_color = RED if is_priority else INK
        text(s, x + Inches(0.95), y + Inches(0.35), card_w - Inches(1.2), Inches(0.5),
             title, size=14, bold=True, color=title_color)
        text(s, x + Inches(0.95), y + Inches(0.85), card_w - Inches(1.2), Inches(1.6),
             body, size=11, color=GREY_TXT)


# === Build ===

def build(out_path):
    prs = prs_setup()
    slide1_cover(prs)
    slide2_problem_vision(prs)
    slide3_objectives(prs)
    slide4_gantt(prs)
    slide5_deliverables(prs)
    slide6_unige(prs)
    slide7_buccarella(prs)
    slide8_integrated(prs)
    slide9_raci(prs)
    slide10_strategic(prs)
    slide11_next_steps(prs)
    prs.save(str(out_path))
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    out = Path(r"\\wsl.localhost\Ubuntu-20.04\home\ev04\v2\ARISE_Kickoff_Meeting_Corrected.pptx")
    build(out)
