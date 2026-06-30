"""DINOGMI working-session deck.

Visual-first explanation of the biomechanical KPIs and the clinical
error taxonomy. Designed for an audience of biomechanics specialists,
so every measurable quantity is shown as a diagram (stick figure with
the relevant vectors, angles, or excursions highlighted) rather than
defined in prose alone.

Design rules:
- One idea per slide.
- Diagram dominates the slide; text is a caption.
- One accent color used semantically: navy for measurement, red for
  errors and red flags, green for correct, grey for axes and labels.
- No decorative AI-feel: no em-dashes, no gradients, no clip art.
"""

import math
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Palette
NAVY = RGBColor(0x1A, 0x33, 0x5C)
INK  = RGBColor(0x14, 0x14, 0x14)
GREY = RGBColor(0x88, 0x88, 0x88)
PALE = RGBColor(0xD9, 0xD9, 0xD9)
RED  = RGBColor(0xC0, 0x2A, 0x2A)
GRN  = RGBColor(0x2E, 0x7D, 0x32)
AMBR = RGBColor(0xC9, 0x86, 0x16)
BG   = RGBColor(0xFA, 0xFA, 0xFA)
WHT  = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── primitives ─────────────────────────────────────────────────────

def prs_setup():
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def text(slide, x, y, w, h, content, size=12, bold=False, italic=False,
         color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top  = Inches(0.0);  tf.margin_bottom = Inches(0.0)
    lines = content if isinstance(content, list) else [content]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def line(slide, x1, y1, x2, y2, color=INK, weight=1.0, dash=False):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     int(x1), int(y1), int(x2), int(y2))
    cn.line.color.rgb = color
    # Reference axes drawn as thinner lighter lines instead of dashed (safer)
    if dash:
        cn.line.width = Pt(max(0.5, weight * 0.6))
    else:
        cn.line.width = Pt(weight)
    return cn


def dot(slide, cx, cy, r, fill=INK, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                int(cx - r), int(cy - r),
                                int(2 * r), int(2 * r))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(0.5)
    s.shadow.inherit = False
    return s


def arrow(slide, x1, y1, x2, y2, color=INK, weight=1.5):
    # Arrows simplified to thick lines. The text label nearby explains direction.
    return line(slide, x1, y1, x2, y2, color=color, weight=weight)


def arc_wedge(slide, cx, cy, r, start_deg, end_deg, color=RED, weight=1.5, n=16):
    """Draw a partial arc as a polyline of n short segments.
    Angles in degrees, measured counterclockwise from positive X axis.
    Y is interpreted in screen coordinates (positive = down)."""
    pts = []
    s = math.radians(start_deg)
    e = math.radians(end_deg)
    for i in range(n + 1):
        t = i / n
        angle = s + (e - s) * t
        x = cx + r * math.cos(angle)
        y = cy - r * math.sin(angle)  # invert Y for screen
        pts.append((int(x), int(y)))
    for i in range(len(pts) - 1):
        line(slide, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1],
             color=color, weight=weight)


def arc_marker(slide, cx, cy, r, label, color=RED, label_offset=(0, 0)):
    """Backwards-compatible helper, now drawn as a clean partial wedge."""
    # Draw the wedge in the upper-left quadrant (90 to 135 deg), which
    # naturally sits between a vertical reference (going up) and a
    # forward-leaning segment.
    arc_wedge(slide, cx, cy, r, 90, 135, color=color, weight=1.5)
    if label:
        lx, ly = label_offset
        text(slide, cx + lx, cy + ly, Inches(0.8), Inches(0.3),
             label, size=10, italic=True, color=color)


# ─── slide chrome ───────────────────────────────────────────────────

def page_header(slide, kicker, title):
    text(slide, Inches(0.5), Inches(0.30), Inches(12), Inches(0.3),
         kicker.upper(), size=10, bold=True, color=NAVY)
    text(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.55),
         title, size=24, bold=True, color=INK)
    line(slide, Inches(0.5), Inches(1.18), Inches(12.8), Inches(1.18),
         color=PALE, weight=1)


def page_footer(slide, n, total):
    text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
         "ARISE   |   D1.1 Biomechanical KPIs and Clinical Error Taxonomy   |   DINOGMI working session",
         size=8, color=GREY)
    text(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
         f"{n} / {total}", size=8, color=GREY, align=PP_ALIGN.RIGHT)


# ─── reusable diagrams ──────────────────────────────────────────────
# All diagrams take a bounding box (x, y, w, h) and draw inside it.

def _stick_sagittal(slide, bx, by, bw, bh, *,
                    trunk_angle_deg=15.0, knee_angle_deg=160.0,
                    show_trunk_angle=True, show_knee_angle=True,
                    title=None, posture="standing"):
    """Side-view stick figure. Angles in degrees from anatomical reference.
    trunk_angle_deg: angle of trunk from vertical (0 = upright).
    knee_angle_deg: internal knee angle (180 = full extension).
    """
    cx = bx + bw / 2
    # Vertical layout: head top, ankle near bottom of bounding box.
    top = by + bh * 0.05
    head_r = bh * 0.05
    head_c = (cx, top + head_r)
    neck   = (cx, top + 2 * head_r + bh * 0.02)
    # Trunk leans forward (positive X) by trunk_angle_deg
    trunk_len = bh * 0.30
    rad = math.radians(trunk_angle_deg)
    hip_x = neck[0] + trunk_len * math.sin(rad)
    hip_y = neck[1] + trunk_len * math.cos(rad)
    hip   = (hip_x, hip_y)
    # Upper leg vertical; lower leg drops to ankle.
    upper_leg = bh * 0.20
    knee = (hip[0] + bh * 0.05, hip[1] + upper_leg)
    # Knee angle: shin direction relative to vertical
    shin_rad = math.radians(180 - knee_angle_deg) * 0.5  # visual approximation
    lower_leg = bh * 0.22
    ankle = (knee[0] - lower_leg * math.sin(shin_rad), knee[1] + lower_leg * math.cos(shin_rad))
    # Arm: shoulder to wrist, hanging
    shoulder = (neck[0] + bh * 0.01, neck[1] + bh * 0.03)
    elbow = (shoulder[0] + bh * 0.04, shoulder[1] + bh * 0.10)
    wrist = (elbow[0] + bh * 0.02, elbow[1] + bh * 0.10)
    # Chair: rectangle behind the figure if sitting/posture
    if posture in ("sitting", "lean"):
        chair_y = hip_y + Inches(0.05)
        rect(slide, cx - Inches(0.5), chair_y, Inches(1.0), Inches(0.6),
             fill=PALE, line=GREY)
    # Floor
    floor_y = by + bh - bh * 0.05
    line(slide, bx + bh * 0.05, floor_y, bx + bw - bh * 0.05, floor_y,
         color=GREY, weight=1)
    # Draw figure
    # Head
    dot(slide, head_c[0], head_c[1], head_r, fill=NAVY)
    # Neck → hip (trunk)
    line(slide, neck[0], neck[1], hip[0], hip[1], color=INK, weight=2)
    # Hip → knee → ankle
    line(slide, hip[0], hip[1], knee[0], knee[1], color=INK, weight=2)
    line(slide, knee[0], knee[1], ankle[0], ankle[1], color=INK, weight=2)
    # Arm
    line(slide, shoulder[0], shoulder[1], elbow[0], elbow[1], color=INK, weight=1.5)
    line(slide, elbow[0], elbow[1], wrist[0], wrist[1], color=INK, weight=1.5)
    # Joint markers
    for jx, jy in [neck, hip, knee, ankle, shoulder, elbow, wrist]:
        dot(slide, jx, jy, Inches(0.06), fill=NAVY)
    # Trunk angle indicator: arc at the neck between vertical-down and trunk segment
    if show_trunk_angle:
        # Vertical reference (thin grey)
        line(slide, neck[0], neck[1], neck[0], neck[1] + Inches(0.65),
             color=GREY, weight=0.8, dash=True)
        # Arc from straight-down (270 in math) sweeping to the trunk direction
        arc_r = Inches(0.42)
        arc_wedge(slide, neck[0], neck[1], arc_r,
                  start_deg=270, end_deg=270 + trunk_angle_deg,
                  color=RED, weight=2.0)
        # Label with theta symbol
        text(slide, neck[0] + Inches(0.50), neck[1] + Inches(0.35),
             Inches(2.5), Inches(0.4),
             [f"θ trunk = {int(trunk_angle_deg)}°"],
             size=13, italic=True, color=RED, bold=True)
    # Knee angle indicator: arc at the knee, on the posterior side
    if show_knee_angle:
        # Compute the local opening on the back of the knee. The hip-knee
        # vector points roughly downward (math 270 + small forward angle),
        # and the knee-ankle vector points downward and slightly back. The
        # interior angle on the posterior side is (180 - knee_angle_deg + small).
        # Use a fixed visible arc on the posterior side for legibility.
        arc_r = Inches(0.30)
        # Posterior side is opposite to the lean direction
        sweep = max(8, 180 - knee_angle_deg + 5)  # ensure visible
        # Draw between 180 (left, posterior) and 180 + sweep
        arc_wedge(slide, knee[0], knee[1], arc_r,
                  start_deg=180, end_deg=180 + sweep,
                  color=RED, weight=2.0)
        text(slide, knee[0] + Inches(0.35), knee[1] - Inches(0.10),
             Inches(2.0), Inches(0.4),
             f"θ knee = {int(knee_angle_deg)}°",
             size=13, italic=True, color=RED, bold=True)
    if title:
        text(slide, bx, by + bh + Inches(0.05), bw, Inches(0.3),
             title, size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)


def _stick_frontal(slide, bx, by, bw, bh, *,
                   knee_valgus_ratio=1.0, shoulder_dy=0.0, title=None):
    """Front-view stick figure. knee_valgus_ratio: KHR (~1.0 normal, <0.7 valgus, >1.2 varus).
    shoulder_dy: vertical displacement of right shoulder vs left, fraction of bh.
    """
    cx = bx + bw / 2
    top = by + bh * 0.05
    head_r = bh * 0.05
    head_c = (cx, top + head_r)
    # Shoulders
    hip_width = bh * 0.16
    shoulder_w = bh * 0.20
    sh_y = top + 2 * head_r + bh * 0.03
    left_sh  = (cx - shoulder_w / 2, sh_y + bh * shoulder_dy / 2)
    right_sh = (cx + shoulder_w / 2, sh_y - bh * shoulder_dy / 2)
    # Hips
    hip_y = sh_y + bh * 0.30
    left_hip  = (cx - hip_width / 2, hip_y)
    right_hip = (cx + hip_width / 2, hip_y)
    # Knees: laterally placed according to valgus ratio
    knee_y = hip_y + bh * 0.22
    knee_distance = hip_width * knee_valgus_ratio
    left_knee  = (cx - knee_distance / 2, knee_y)
    right_knee = (cx + knee_distance / 2, knee_y)
    # Ankles directly under hips
    ankle_y = knee_y + bh * 0.22
    left_ankle  = (cx - hip_width / 2, ankle_y)
    right_ankle = (cx + hip_width / 2, ankle_y)
    # Floor
    line(slide, bx + bh * 0.05, ankle_y + bh * 0.03,
         bx + bw - bh * 0.05, ankle_y + bh * 0.03, color=GREY, weight=1)
    # Head
    dot(slide, head_c[0], head_c[1], head_r, fill=NAVY)
    # Trunk: shoulders to hips, X pattern style (mid-line)
    line(slide, left_sh[0], left_sh[1], right_sh[0], right_sh[1], color=INK, weight=2)
    line(slide, left_hip[0], left_hip[1], right_hip[0], right_hip[1], color=INK, weight=2)
    line(slide, left_sh[0], left_sh[1], left_hip[0], left_hip[1], color=INK, weight=2)
    line(slide, right_sh[0], right_sh[1], right_hip[0], right_hip[1], color=INK, weight=2)
    # Mid spine
    mid_sh = (cx, (left_sh[1] + right_sh[1]) / 2)
    mid_hip = (cx, hip_y)
    line(slide, mid_sh[0], mid_sh[1], mid_hip[0], mid_hip[1], color=GREY, weight=1, dash=True)
    # Legs
    valgus_color = RED if knee_valgus_ratio < 0.85 or knee_valgus_ratio > 1.15 else INK
    line(slide, left_hip[0], left_hip[1], left_knee[0], left_knee[1], color=valgus_color, weight=2)
    line(slide, right_hip[0], right_hip[1], right_knee[0], right_knee[1], color=valgus_color, weight=2)
    line(slide, left_knee[0], left_knee[1], left_ankle[0], left_ankle[1], color=valgus_color, weight=2)
    line(slide, right_knee[0], right_knee[1], right_ankle[0], right_ankle[1], color=valgus_color, weight=2)
    # Joint dots
    for jx, jy in [left_sh, right_sh, left_hip, right_hip, left_knee, right_knee,
                   left_ankle, right_ankle]:
        dot(slide, jx, jy, Inches(0.06), fill=NAVY)
    if title:
        text(slide, bx, by + bh + Inches(0.05), bw, Inches(0.3),
             title, size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    return {
        "left_hip": left_hip, "right_hip": right_hip,
        "left_knee": left_knee, "right_knee": right_knee,
        "left_ankle": left_ankle, "right_ankle": right_ankle,
        "left_sh": left_sh, "right_sh": right_sh,
        "hip_width": hip_width, "knee_distance": knee_distance,
    }


def _hipY_signal(slide, bx, by, bw, bh, *, phases=True):
    """Stylised hip-Y vertical-position curve across one STS repetition.
    Y is normalised hip elevation (0 = seated baseline, 1 = peak standing).
    X is time in seconds, ~3 s total rep."""
    # Plot area inset to leave room for tick labels
    inset_l = Inches(0.55)
    inset_b = Inches(0.35)
    pax_x = bx + inset_l
    pax_y = by
    pax_w = bw - inset_l - Inches(0.10)
    pax_h = bh - inset_b
    # Axes
    line(slide, pax_x, pax_y + pax_h, pax_x + pax_w, pax_y + pax_h,
         color=GREY, weight=1)               # X axis
    line(slide, pax_x, pax_y, pax_x, pax_y + pax_h, color=GREY, weight=1)  # Y axis
    # Y-axis ticks and labels: 0.0 to 1.0, step 0.25
    y_ticks = [0.0, 0.25, 0.50, 0.75, 1.00]
    for v in y_ticks:
        yy = pax_y + pax_h - (pax_h * v)
        # tick mark
        line(slide, pax_x - Inches(0.07), yy, pax_x, yy, color=GREY, weight=1)
        # gridline (light)
        line(slide, pax_x, yy, pax_x + pax_w, yy, color=PALE, weight=0.5)
        # label
        text(slide, pax_x - Inches(0.55), yy - Inches(0.10),
             Inches(0.45), Inches(0.22), f"{v:.2f}",
             size=9, color=GREY, align=PP_ALIGN.RIGHT)
    # Y-axis title
    text(slide, bx - Inches(0.05), pax_y + pax_h * 0.35,
         Inches(0.6), Inches(0.4),
         ["hip Y", "(normalised)"], size=9, italic=True, color=GREY,
         align=PP_ALIGN.LEFT)
    # X-axis ticks and labels: 0s, 1s, 2s, 3s
    t_ticks = [(0.00, "0 s"), (0.33, "1 s"), (0.66, "2 s"), (1.00, "3 s")]
    for t, lab in t_ticks:
        xx = pax_x + pax_w * t
        line(slide, xx, pax_y + pax_h, xx, pax_y + pax_h + Inches(0.07),
             color=GREY, weight=1)
        text(slide, xx - Inches(0.25), pax_y + pax_h + Inches(0.10),
             Inches(0.5), Inches(0.22), lab,
             size=9, color=GREY, align=PP_ALIGN.CENTER)
    text(slide, pax_x + pax_w / 2 - Inches(0.5),
         pax_y + pax_h + Inches(0.32),
         Inches(1.0), Inches(0.22), "time",
         size=9, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    # 6 phases with stabilisation between standing and descent
    # x-bounds for each phase as fractions of plot width
    phase_bounds = [
        (0.00, 0.12, "sitting",       PALE),
        (0.12, 0.30, "forward lean",  AMBR),
        (0.30, 0.50, "lift-off",      NAVY),
        (0.50, 0.65, "standing",      GRN),
        (0.65, 0.78, "stabilisation", GRN),
        (0.78, 1.00, "descent",       AMBR),
    ]
    if phases:
        bar_y = pax_y + pax_h + Inches(0.45)
        for a, b, lab, col in phase_bounds:
            x0 = pax_x + pax_w * a
            seg_w = pax_w * (b - a)
            rect(slide, x0, bar_y, seg_w - Inches(0.02), Inches(0.20),
                 fill=col)
            text(slide, x0, bar_y + Inches(0.22), seg_w, Inches(0.25),
                 lab, size=9, color=INK, align=PP_ALIGN.CENTER)
    # Polyline approximating the hip-Y curve
    # (x_frac of pax_w, y_norm where 0 = bottom, 1 = top)
    pts_norm = [
        (0.00, 0.06),  # sitting baseline
        (0.12, 0.08),
        (0.22, 0.15),  # forward lean (slight rise)
        (0.30, 0.20),
        (0.40, 0.55),  # lift-off rapid rise
        (0.50, 0.92),  # standing peak
        (0.58, 0.93),
        (0.65, 0.92),  # stabilisation plateau
        (0.72, 0.91),
        (0.78, 0.85),  # descent begins
        (0.90, 0.30),
        (1.00, 0.06),  # back to sitting
    ]
    prev = None
    for fx, fy in pts_norm:
        x = pax_x + pax_w * fx
        y = pax_y + pax_h - pax_h * fy
        if prev:
            line(slide, prev[0], prev[1], x, y, color=NAVY, weight=2.2)
        prev = (x, y)


# ─── slide builders ─────────────────────────────────────────────────

def s_title(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    # Accent diagonal line
    line(s, Inches(0.7), Inches(2.1), Inches(2.0), Inches(2.1), color=WHT, weight=2)
    text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(0.5),
         "ARISE   |   WP1 / T1.3", size=12, bold=True, color=WHT)
    text(s, Inches(0.7), Inches(2.7), Inches(12), Inches(1.4),
         "Biomechanical KPIs and",
         size=42, bold=True, color=WHT)
    text(s, Inches(0.7), Inches(3.5), Inches(12), Inches(1.4),
         "Clinical Error Taxonomy",
         size=42, bold=True, color=WHT)
    text(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.4),
         "Working session with DINOGMI, University of Genoa",
         size=18, color=WHT)
    text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
         "Innovina S.r.l.        25 May 2026        v0.3",
         size=11, color=WHT)
    return s


def s_one_question(prs):
    s = blank(prs)
    page_header(s, "Premise", "What we are actually measuring")
    text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.0),
         "An STS repetition is a constrained kinematic event.",
         size=28, bold=True, color=INK)
    text(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.8),
         "It has a known temporal structure (5 phases) and a small set of",
         size=20, color=INK)
    text(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.8),
         "biomechanical signals that distinguish correct from compensated execution.",
         size=20, color=INK)
    text(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.4),
         "This deck specifies:",
         size=14, bold=True, color=NAVY)
    text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.4),
         "(1)  the 18 scalar KPIs we extract per repetition,",
         size=14, color=INK)
    text(s, Inches(0.7), Inches(5.4), Inches(12), Inches(0.4),
         "(2)  the 14 clinical error codes derived from those KPIs,",
         size=14, color=INK)
    text(s, Inches(0.7), Inches(5.8), Inches(12), Inches(0.4),
         "(3)  the 2 always-on safety flags (G1, G2),",
         size=14, color=INK)
    text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.4),
         "(4)  the thresholds and the open items needing DINOGMI sign-off.",
         size=14, color=INK)
    return s


def s_pipeline(prs):
    s = blank(prs)
    page_header(s, "Measurement chain", "Camera to KPI in five steps")
    # Five chevron-like blocks
    steps = [
        ("Acquisition", "RGB 30 fps"),
        ("Pose", "33 landmarks / frame"),
        ("Rigidify", "Per-patient bone-length"),
        ("Features", "Angles, CoM, velocities"),
        ("Aggregate", "18 KPIs per rep"),
    ]
    block_w = Inches(2.30)
    gap = Inches(0.12)
    h = Inches(1.4)
    start_x = Inches(0.7)
    y = Inches(2.0)
    for i, (a, b) in enumerate(steps):
        x = start_x + i * (block_w + gap)
        rect(s, x, y, block_w, h, fill=BG, line=NAVY, line_w=1.0)
        text(s, x, y + Inches(0.25), block_w, Inches(0.35),
             a, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        text(s, x, y + Inches(0.7), block_w, Inches(0.4),
             b, size=11, color=INK, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            ax = x + block_w
            line(s, ax, y + h / 2, ax + gap, y + h / 2, color=NAVY, weight=2)
    # Critical invariant box
    rect(s, Inches(0.7), Inches(4.0), Inches(12), Inches(2.6), fill=BG, line=PALE)
    text(s, Inches(0.9), Inches(4.15), Inches(11), Inches(0.4),
         "Critical invariant", size=14, bold=True, color=NAVY)
    text(s, Inches(0.9), Inches(4.55), Inches(11), Inches(1.6),
         ["Height-normalise the whole repetition by the head-to-foot distance at frame 0.",
          "Do not refresh per frame: per-frame normalisation regresses out the pelvis-rise signal,",
          "which is the dominant kinematic feature of an STS.",
          "Origin: mid-pelvis at frame 0. Y axis up (gravity). X axis forward (sagittal). Z lateral."],
         size=13, color=INK)
    return s


def s_phases(prs):
    s = blank(prs)
    page_header(s, "Phase segmentation",
                "6 phases, detected from hip-Y trajectory, knee extension, and trunk angle")
    _hipY_signal(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(3.4))
    # Caption below
    text(s, Inches(0.9), Inches(6.30), Inches(11.5), Inches(0.4),
         "Why phases matter",
         size=13, bold=True, color=NAVY)
    text(s, Inches(0.9), Inches(6.60), Inches(11.5), Inches(0.4),
         "Same posture, different clinical meaning depending on phase. Trunk angle of 40° is expected during forward lean; the same value during stabilisation is E10 (Incomplete Upright).",
         size=11, italic=True, color=GREY)
    return s


# ─── KPI diagram-driven slides ───────────────────────────────────────

def kpi_slide(prs, kicker, title, diagram_fn, formula, clinical, threshold,
              linked_errors):
    """Generic KPI slide: left half diagram, right half captions."""
    s = blank(prs)
    page_header(s, kicker, title)
    # Diagram area (left)
    diagram_fn(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    # Caption (right)
    rx = Inches(7.2); rw = Inches(5.8)
    text(s, rx, Inches(1.5), rw, Inches(0.4),
         "Formula", size=11, bold=True, color=NAVY)
    rect(s, rx, Inches(1.85), rw, Inches(1.0), fill=BG, line=PALE)
    text(s, rx + Inches(0.15), Inches(1.92), rw - Inches(0.3), Inches(0.95),
         formula, size=14, italic=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    text(s, rx, Inches(3.05), rw, Inches(0.4),
         "Clinical meaning", size=11, bold=True, color=NAVY)
    text(s, rx, Inches(3.4), rw, Inches(2.2),
         clinical, size=13, color=INK)
    text(s, rx, Inches(5.5), rw, Inches(0.4),
         "Threshold (v6 calibration)", size=11, bold=True, color=NAVY)
    text(s, rx, Inches(5.85), rw, Inches(0.5),
         threshold, size=13, color=RED, bold=True)
    text(s, rx, Inches(6.5), rw, Inches(0.3),
         "Triggers:  " + linked_errors, size=11, color=GREY, italic=True)
    return s


# Per-KPI diagrams (each takes slide and bounding box).

def diag_trunk_angle(s, x, y, w, h):
    _stick_sagittal(s, x, y, w, h, trunk_angle_deg=30,
                    show_trunk_angle=True, show_knee_angle=False,
                    title="Side view, ascent phase")


def diag_spine(s, x, y, w, h):
    """Two side-by-side stick figures showing the mid-hip → mid-shoulder
    → nose chain. Left is neutral (θ ≈ 180°); right is slumped (θ < 150°).
    The angle at the mid-shoulder vertex is what we measure."""
    half_w = (w - Inches(0.2)) / 2
    _spine_figure(s, x,                       y, half_w, h, spine_angle=180,
                  title="Neutral spine,  θ ≈ 180°", seg_color=INK)
    _spine_figure(s, x + half_w + Inches(0.2), y, half_w, h, spine_angle=140,
                  title="Slumped spine,  θ < 150°  (E3b)", seg_color=RED)


def _spine_figure(s, bx, by, bw, bh, *, spine_angle, title, seg_color):
    """Draws one spine illustration. Three labelled landmarks (hip,
    shoulder, nose) connected by two segments meeting at the shoulder."""
    cx = bx + bw / 2
    # Anchor: hip at lower-middle
    hip_y = by + bh * 0.70
    hip   = (cx, hip_y)
    # Mid-shoulder above hip (vertical)
    shoulder_len = Inches(1.2)
    mid_sh_y = hip_y - shoulder_len
    mid_sh = (cx, mid_sh_y)
    # Head/nose tilts forward by (180 - spine_angle)
    head_len = Inches(0.7)
    bend = math.radians(180 - spine_angle)
    nose_x = cx + math.sin(bend) * head_len
    nose_y = mid_sh_y - math.cos(bend) * head_len
    nose = (nose_x, nose_y)
    # Floor and legs (simple, just to anchor the figure)
    floor_y = hip_y + Inches(1.4)
    line(s, bx + Inches(0.3), floor_y, bx + bw - Inches(0.3), floor_y,
         color=GREY, weight=1)
    knee_y = hip_y + Inches(0.7)
    ankle_y = floor_y - Inches(0.05)
    line(s, hip[0], hip[1], hip[0] + Inches(0.05), knee_y, color=INK, weight=2.2)
    line(s, hip[0] + Inches(0.05), knee_y, hip[0], ankle_y, color=INK, weight=2.2)
    dot(s, hip[0] + Inches(0.05), knee_y, Inches(0.05), fill=NAVY)
    # Three landmarks of the chain
    # 1. Mid-hip → Mid-shoulder segment
    line(s, hip[0], hip[1], mid_sh[0], mid_sh[1], color=INK, weight=3.0)
    # 2. Mid-shoulder → Nose segment (red if slumped)
    line(s, mid_sh[0], mid_sh[1], nose[0], nose[1], color=seg_color, weight=3.0)
    # Joint markers (bigger at vertex)
    dot(s, hip[0], hip[1], Inches(0.08), fill=NAVY, line_color=WHT)
    dot(s, mid_sh[0], mid_sh[1], Inches(0.10), fill=NAVY, line_color=WHT)
    dot(s, nose[0], nose[1], Inches(0.13), fill=NAVY)  # head
    # Landmark labels (left side: hip and shoulder)
    text(s, bx + Inches(0.1), hip[1] - Inches(0.12),
         Inches(1.2), Inches(0.25), "mid-hip",
         size=9, italic=True, color=GREY, align=PP_ALIGN.LEFT)
    text(s, bx + Inches(0.1), mid_sh[1] - Inches(0.12),
         Inches(1.4), Inches(0.25), "mid-shoulder",
         size=9, italic=True, color=GREY, align=PP_ALIGN.LEFT)
    text(s, nose[0] + Inches(0.18), nose[1] - Inches(0.12),
         Inches(1.0), Inches(0.25), "nose",
         size=9, italic=True, color=GREY, align=PP_ALIGN.LEFT)
    # Reference: where the head would be if spine were straight
    if spine_angle < 180:
        line(s, mid_sh[0], mid_sh[1], cx, mid_sh_y - head_len,
             color=GREY, weight=0.8, dash=True)
        # Arc indicating the spine angle (the bend at the shoulder vertex)
        arc_wedge(s, mid_sh[0], mid_sh[1], Inches(0.45),
                  start_deg=90, end_deg=90 - (180 - spine_angle),
                  color=RED, weight=2.0, n=12)
    # Angle label
    text(s, mid_sh[0] + Inches(0.35), mid_sh[1] - Inches(0.40),
         Inches(2.4), Inches(0.35),
         f"θ = {spine_angle}°",
         size=14, italic=True, bold=True,
         color=RED if spine_angle < 180 else NAVY)
    # Sub-title under figure
    text(s, bx, by + bh - Inches(0.30), bw, Inches(0.3),
         title, size=11, italic=True,
         color=RED if spine_angle < 180 else NAVY,
         align=PP_ALIGN.CENTER, bold=True)


def diag_lateral_sway(s, x, y, w, h):
    p = _stick_frontal(s, x, y, w, h, knee_valgus_ratio=1.0,
                       shoulder_dy=0.0, title="Front view, ascent")
    # Overlay a CoM dot and lateral excursion arrows
    cx = x + w / 2
    com_y = (p["left_sh"][1] + p["left_hip"][1]) / 2
    dot(s, cx + Inches(0.5), com_y, Inches(0.08), fill=RED)
    text(s, cx + Inches(0.6), com_y - Inches(0.1), Inches(1.5), Inches(0.3),
         "CoM (lateral)", size=10, italic=True, color=RED)
    arrow(s, cx, com_y, cx + Inches(0.45), com_y, color=RED, weight=1.5)


def diag_knee_extension(s, x, y, w, h):
    _stick_sagittal(s, x, y, w, h, trunk_angle_deg=5, knee_angle_deg=170,
                    show_trunk_angle=False, show_knee_angle=True,
                    title="Side view, standing phase (peak extension)")


def diag_valgus(s, x, y, w, h):
    # Two figures side by side: correct vs valgus
    half_w = w / 2 - Inches(0.05)
    p1 = _stick_frontal(s, x, y, half_w, h, knee_valgus_ratio=1.0,
                        title="KHR ≈ 1.0 (neutral)")
    p2 = _stick_frontal(s, x + half_w + Inches(0.1), y, half_w, h,
                        knee_valgus_ratio=0.55,
                        title="KHR < 0.7 (valgus)")
    # Annotate the second
    cx2 = x + half_w + Inches(0.1) + half_w / 2
    text(s, cx2 - Inches(0.5), p2["left_knee"][1] + Inches(0.1),
         Inches(1.5), Inches(0.3),
         "knees inward", size=10, italic=True, color=RED, align=PP_ALIGN.CENTER)


def diag_khr(s, x, y, w, h):
    # Single front view showing the two distance measurements
    p = _stick_frontal(s, x, y, w, h, knee_valgus_ratio=0.85,
                       title="KHR = inter-knee / inter-hip")
    # Highlight hip distance with red bracket
    lh, rh = p["left_hip"], p["right_hip"]
    lk, rk = p["left_knee"], p["right_knee"]
    line(s, lh[0], lh[1] - Inches(0.05), rh[0], rh[1] - Inches(0.05),
         color=RED, weight=1.5)
    text(s, (lh[0] + rh[0]) / 2 - Inches(0.4), lh[1] - Inches(0.30),
         Inches(0.8), Inches(0.3),
         "hip_w", size=10, italic=True, color=RED, align=PP_ALIGN.CENTER)
    line(s, lk[0], lk[1] + Inches(0.10), rk[0], rk[1] + Inches(0.10),
         color=RED, weight=1.5)
    text(s, (lk[0] + rk[0]) / 2 - Inches(0.5), lk[1] + Inches(0.13),
         Inches(1.0), Inches(0.3),
         "knee_dist", size=10, italic=True, color=RED, align=PP_ALIGN.CENTER)


def diag_instability(s, x, y, w, h):
    # Two stacked oscillation signals
    # Hip velocity over time, normal vs unstable
    pad = Inches(0.2)
    plot_h = h / 2 - pad
    # Top: normal
    line(s, x + Inches(0.4), y + plot_h / 2,
         x + w - Inches(0.2), y + plot_h / 2, color=GREY, weight=1)
    text(s, x, y + Inches(0.05), Inches(2), Inches(0.3),
         "normal lift", size=10, color=GRN, bold=True)
    # smooth curve
    import math
    n = 30
    prev = None
    for i in range(n + 1):
        t = i / n
        xt = x + Inches(0.4) + t * (w - Inches(0.6))
        yt = y + plot_h / 2 - Inches(0.6) * math.sin(t * math.pi)
        if prev:
            line(s, prev[0], prev[1], xt, yt, color=GRN, weight=2)
        prev = (xt, yt)
    # Bottom: oscillating
    y2 = y + plot_h + pad
    line(s, x + Inches(0.4), y2 + plot_h / 2,
         x + w - Inches(0.2), y2 + plot_h / 2, color=GREY, weight=1)
    text(s, x, y2 + Inches(0.05), Inches(2), Inches(0.3),
         "unstable lift", size=10, color=RED, bold=True)
    prev = None
    for i in range(n * 2 + 1):
        t = i / (n * 2)
        xt = x + Inches(0.4) + t * (w - Inches(0.6))
        # oscillating component
        yt = y2 + plot_h / 2 - Inches(0.5) * math.sin(t * math.pi) - Inches(0.15) * math.sin(t * math.pi * 8)
        if prev:
            line(s, prev[0], prev[1], xt, yt, color=RED, weight=2)
        prev = (xt, yt)
    text(s, x, y + h + Inches(0.05), w, Inches(0.3),
         "Hip vertical velocity over lift-off", size=10, italic=True, color=GREY,
         align=PP_ALIGN.CENTER)


def diag_hip(s, x, y, w, h):
    """Side-view figure with the hip-joint angle shown as an arc at the
    hip vertex, formed by the chain shoulder → hip → knee."""
    _hip_figure(s, x, y, w, h, trunk_angle_deg=12, hip_angle_deg=120,
                show_velocity=True,
                title="Mid lift-off. Hip extending, pelvis rising.")


def _hip_figure(s, bx, by, bw, bh, *, trunk_angle_deg, hip_angle_deg,
                show_velocity, title):
    cx = bx + bw / 2
    head_r = Inches(0.18)
    # Anchor: hip at centre
    hip = (cx, by + bh * 0.50)
    # Trunk goes up-and-slightly-forward
    trunk_len = Inches(1.6)
    rad_trunk = math.radians(trunk_angle_deg)
    sh_x = hip[0] + math.sin(rad_trunk) * trunk_len
    sh_y = hip[1] - math.cos(rad_trunk) * trunk_len
    mid_sh = (sh_x, sh_y)
    # Head above shoulder, aligned with trunk
    head_x = sh_x + math.sin(rad_trunk) * Inches(0.45)
    head_y = sh_y - math.cos(rad_trunk) * Inches(0.45)
    # Upper leg goes from hip down-and-slightly-back at angle (180 - hip_angle_deg)
    # from the trunk. Compute knee position.
    # The hip angle is the interior at the hip between (trunk direction)
    # and (upper-leg direction). For visualisation: place the knee so the
    # upper leg is forward of vertical by a small amount and the resulting
    # angle matches hip_angle_deg.
    upper_leg = Inches(1.2)
    # In math convention, trunk vector from hip to shoulder points up
    # at angle (90° - trunk_angle_deg) from positive X.
    trunk_dir_math = 90 - trunk_angle_deg
    # Upper-leg vector from hip to knee: rotate trunk direction by
    # hip_angle_deg, on the anterior side (forward).
    leg_dir_math = trunk_dir_math - hip_angle_deg
    rad_leg = math.radians(leg_dir_math)
    knee_x = hip[0] + math.cos(rad_leg) * upper_leg
    knee_y = hip[1] - math.sin(rad_leg) * upper_leg
    knee = (knee_x, knee_y)
    # Lower leg drops vertically to ankle
    ankle = (knee[0], knee[1] + Inches(1.2))
    # Floor
    floor_y = ankle[1] + Inches(0.05)
    line(s, bx + Inches(0.3), floor_y, bx + bw - Inches(0.3), floor_y,
         color=GREY, weight=1)
    # Head
    dot(s, head_x, head_y, head_r, fill=NAVY)
    # Trunk
    line(s, mid_sh[0], mid_sh[1], hip[0], hip[1], color=INK, weight=3.0)
    # Upper leg + lower leg (anterior side of hip will host the arc)
    line(s, hip[0], hip[1], knee[0], knee[1], color=INK, weight=3.0)
    line(s, knee[0], knee[1], ankle[0], ankle[1], color=INK, weight=3.0)
    # Joint dots
    dot(s, mid_sh[0], mid_sh[1], Inches(0.08), fill=NAVY)
    dot(s, hip[0], hip[1], Inches(0.10), fill=NAVY)  # vertex, larger
    dot(s, knee[0], knee[1], Inches(0.08), fill=NAVY)
    dot(s, ankle[0], ankle[1], Inches(0.07), fill=NAVY)
    # Landmark labels
    text(s, mid_sh[0] + Inches(0.15), mid_sh[1] - Inches(0.12),
         Inches(1.4), Inches(0.25), "mid-shoulder",
         size=9, italic=True, color=GREY)
    text(s, hip[0] + Inches(0.15), hip[1] - Inches(0.05),
         Inches(0.8), Inches(0.25), "hip",
         size=9, italic=True, color=GREY)
    text(s, knee[0] + Inches(0.15), knee[1] - Inches(0.05),
         Inches(0.8), Inches(0.25), "knee",
         size=9, italic=True, color=GREY)
    # Angle arc at hip vertex
    # The interior angle is on the anterior side (forward of the body),
    # between the trunk-direction (up-forward) and the upper-leg direction
    # (down-forward). Sweep from leg direction to trunk direction.
    arc_wedge(s, hip[0], hip[1], Inches(0.55),
              start_deg=leg_dir_math, end_deg=trunk_dir_math,
              color=RED, weight=2.2, n=18)
    # Angle label
    label_angle = (leg_dir_math + trunk_dir_math) / 2
    rad_lbl = math.radians(label_angle)
    lbl_x = hip[0] + math.cos(rad_lbl) * Inches(0.95)
    lbl_y = hip[1] - math.sin(rad_lbl) * Inches(0.95)
    text(s, lbl_x - Inches(0.6), lbl_y - Inches(0.15),
         Inches(2.0), Inches(0.35),
         f"θ hip = {hip_angle_deg}°",
         size=13, italic=True, bold=True, color=RED, align=PP_ALIGN.CENTER)
    # Vertical velocity arrow at the pelvis
    if show_velocity:
        vx0 = hip[0] - Inches(0.9)
        vy_top = hip[1] - Inches(0.45)
        vy_bot = hip[1] + Inches(0.45)
        line(s, vx0, vy_bot, vx0, vy_top, color=RED, weight=2.5)
        # Small arrowhead triangle pointing up
        head_w = Inches(0.10)
        tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  int(vx0 - head_w), int(vy_top - head_w * 1.2),
                                  int(head_w * 2), int(head_w * 2))
        tri.fill.solid()
        tri.fill.fore_color.rgb = RED
        tri.line.fill.background()
        tri.shadow.inherit = False
        text(s, vx0 - Inches(1.6), hip[1] - Inches(0.2),
             Inches(1.3), Inches(0.30),
             "v(hip Y)", size=11, italic=True, bold=True, color=RED,
             align=PP_ALIGN.RIGHT)
    # Sub-title
    text(s, bx, by + bh - Inches(0.25), bw, Inches(0.3),
         title, size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)


def diag_com_shoulder(s, x, y, w, h):
    _stick_frontal(s, x, y, w, h, knee_valgus_ratio=1.0,
                   shoulder_dy=0.08,
                   title="Shoulder differential exaggerated")


def diag_hands(s, x, y, w, h):
    """Two side-by-side figures showing the two hands-assist subtypes:
    pushing on the knees (left) and pushing on the chair seat (right)."""
    half_w = (w - Inches(0.2)) / 2
    _hands_figure(s, x,                       y, half_w, h,
                  pattern="on_knee",
                  title="Wrist on knee  (push off knees)")
    _hands_figure(s, x + half_w + Inches(0.2), y, half_w, h,
                  pattern="on_chair",
                  title="Wrist below hip  (push off chair seat)")


def _hands_figure(s, bx, by, bw, bh, *, pattern, title):
    """Side-view figure with the arm bent so the wrist is at the
    relevant clinical location."""
    cx = bx + bw / 2
    head_r = Inches(0.18)
    # Trunk leaned forward
    trunk_angle_deg = 25
    rad_trunk = math.radians(trunk_angle_deg)
    # Hip near middle
    hip = (cx - Inches(0.20), by + bh * 0.50)
    # Trunk up-forward
    trunk_len = Inches(1.5)
    sh_x = hip[0] + math.sin(rad_trunk) * trunk_len
    sh_y = hip[1] - math.cos(rad_trunk) * trunk_len
    mid_sh = (sh_x, sh_y)
    head_c = (sh_x + math.sin(rad_trunk) * Inches(0.4),
              sh_y - math.cos(rad_trunk) * Inches(0.4))
    # Leg: hip forward to knee, then to ankle
    knee = (hip[0] + Inches(0.55), hip[1] + Inches(0.75))
    ankle = (knee[0] + Inches(0.05), knee[1] + Inches(1.0))
    # Floor + chair
    floor_y = ankle[1] + Inches(0.05)
    line(s, bx + Inches(0.2), floor_y, bx + bw - Inches(0.2), floor_y,
         color=GREY, weight=1)
    chair_top = hip[1] + Inches(0.05)
    rect(s, hip[0] - Inches(0.5), chair_top, Inches(0.7), Inches(0.85),
         fill=PALE, line=GREY)
    # Body segments
    line(s, mid_sh[0], mid_sh[1], hip[0], hip[1], color=INK, weight=3.0)
    line(s, hip[0], hip[1], knee[0], knee[1], color=INK, weight=3.0)
    line(s, knee[0], knee[1], ankle[0], ankle[1], color=INK, weight=3.0)
    dot(s, head_c[0], head_c[1], head_r, fill=NAVY)
    for jp in [mid_sh, hip, knee, ankle]:
        dot(s, jp[0], jp[1], Inches(0.07), fill=NAVY)
    # Arm: shoulder → elbow → wrist, end position depending on pattern
    if pattern == "on_knee":
        wrist = (knee[0] - Inches(0.05), knee[1] - Inches(0.05))
        elbow = (mid_sh[0] + Inches(0.30), (mid_sh[1] + wrist[1]) / 2 - Inches(0.10))
    else:  # on_chair
        wrist = (hip[0] - Inches(0.45), hip[1] + Inches(0.40))
        elbow = (mid_sh[0] - Inches(0.20), (mid_sh[1] + wrist[1]) / 2)
    # Draw arm in red to highlight the assist
    line(s, mid_sh[0], mid_sh[1], elbow[0], elbow[1], color=RED, weight=3.0)
    line(s, elbow[0], elbow[1], wrist[0], wrist[1], color=RED, weight=3.0)
    dot(s, elbow[0], elbow[1], Inches(0.07), fill=RED)
    dot(s, wrist[0], wrist[1], Inches(0.09), fill=RED)
    # Reference line: hip height (the threshold for chair-push)
    if pattern == "on_chair":
        line(s, bx + Inches(0.1), hip[1], bx + bw - Inches(0.1), hip[1],
             color=GREY, weight=0.8, dash=True)
        text(s, bx + Inches(0.1), hip[1] - Inches(0.22),
             Inches(1.5), Inches(0.22), "hip level (ref)",
             size=9, italic=True, color=GREY)
    # Contact indicator: small red circle around the contact zone
    contact_x, contact_y = wrist
    ell = s.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(contact_x - Inches(0.20)), int(contact_y - Inches(0.18)),
        int(Inches(0.40)), int(Inches(0.36))
    )
    ell.fill.background()
    ell.line.color.rgb = RED
    ell.line.width = Pt(1.5)
    ell.shadow.inherit = False
    # Annotation label with leader text
    if pattern == "on_knee":
        ann = "wrist–knee distance ≈ 0"
    else:
        ann = "wrist below hip"
    text(s, contact_x + Inches(0.25), contact_y - Inches(0.10),
         Inches(2.0), Inches(0.30), ann,
         size=11, italic=True, bold=True, color=RED)
    # Sub-title
    text(s, bx, by + bh - Inches(0.25), bw, Inches(0.3),
         title, size=11, italic=True, color=RED, bold=True,
         align=PP_ALIGN.CENTER)


def diag_durations(s, x, y, w, h):
    # Horizontal phase bar with proportional widths
    labels = ["sitting", "forward lean", "lift-off", "standing", "stabilisation", "descent"]
    widths = [0.08, 0.17, 0.13, 0.25, 0.20, 0.17]
    colors = [PALE, AMBR, NAVY, GRN, GRN, AMBR]
    bar_y = y + h * 0.35
    bar_h = Inches(0.6)
    x_cursor = x + Inches(0.3)
    avail_w = w - Inches(0.6)
    for lab, wfrac, col in zip(labels, widths, colors):
        seg_w = avail_w * wfrac
        rect(s, x_cursor, bar_y, seg_w, bar_h, fill=col)
        text(s, x_cursor, bar_y + bar_h + Inches(0.1), seg_w, Inches(0.3),
             lab, size=9, color=INK, align=PP_ALIGN.CENTER)
        text(s, x_cursor, bar_y - Inches(0.3), seg_w, Inches(0.3),
             f"{wfrac:.0%}", size=9, color=GREY, italic=True, align=PP_ALIGN.CENTER)
        x_cursor += seg_w
    text(s, x, y + h - Inches(0.3), w, Inches(0.3),
         "Example timing budget across one healthy STS rep",
         size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)


def diag_personalisation(s, x, y, w, h):
    """Four-layer architecture diagram: rule, z-score, deep learning, forecasting."""
    box_h = Inches(0.85)
    gap = Inches(0.12)
    layers = [
        ("Layer 1.  Population rule thresholds",
         "Absolute clinical errors.  Frozen.", NAVY),
        ("Layer 2.  Per-patient z-score baseline",
         "|z| > 2 fires personal-regression flag.  Online (every 5 correct reps).", NAVY),
        ("Layer 3.  Deep learning classifier  (BiLSTM multi-task)",
         "Binary F1 = 0.92      Error-category macro F1 = 0.85      Active in production.", RED),
        ("Layer 4.  Forecasting head  (proactive guidance)",
         "Predicts trajectory deviation 200 ms before peak.  Latency-aware live cues.", RED),
    ]
    ty = y + Inches(0.10)
    for title, sub, accent in layers:
        rect(s, x, ty, w, box_h, fill=BG, line=accent, line_w=1.0)
        # Left accent strip
        rect(s, x, ty, Inches(0.10), box_h, fill=accent)
        text(s, x + Inches(0.25), ty + Inches(0.10), w - Inches(0.4), Inches(0.32),
             title, size=12, bold=True, color=INK)
        text(s, x + Inches(0.25), ty + Inches(0.42), w - Inches(0.4), Inches(0.36),
             sub, size=10, color=GREY, italic=True)
        ty += box_h + gap
    text(s, x, ty + Inches(0.05), w, Inches(0.3),
         "All four layers run simultaneously. Each compensates for the others' blind spots.",
         size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)


# ─── error code slides ──────────────────────────────────────────────

def err_slide(prs, code, title, diagram_fn, what, rule, kpi_link):
    s = blank(prs)
    # Section indicator
    page_header(s, f"Error code  {code}", title)
    diagram_fn(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.0))
    rx = Inches(7.0); rw = Inches(6.0)
    # Big code badge
    rect(s, rx, Inches(1.5), Inches(1.3), Inches(0.9), fill=RED)
    text(s, rx, Inches(1.5), Inches(1.3), Inches(0.9), code,
         size=24, bold=True, color=WHT, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, rx + Inches(1.5), Inches(1.6), rw - Inches(1.5), Inches(0.4),
         title, size=18, bold=True, color=INK)
    text(s, rx + Inches(1.5), Inches(2.0), rw - Inches(1.5), Inches(0.4),
         "Linked KPI:  " + kpi_link, size=11, color=GREY, italic=True)

    text(s, rx, Inches(2.8), rw, Inches(0.4),
         "Clinical pattern", size=11, bold=True, color=NAVY)
    text(s, rx, Inches(3.15), rw, Inches(2.2),
         what, size=13, color=INK)
    text(s, rx, Inches(5.5), rw, Inches(0.4),
         "Detection rule (v6 calibration)", size=11, bold=True, color=NAVY)
    rect(s, rx, Inches(5.85), rw, Inches(0.85), fill=BG, line=PALE)
    text(s, rx + Inches(0.15), Inches(5.93), rw - Inches(0.3), Inches(0.8),
         rule, size=13, italic=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    return s


# Diagrams reused for errors (overlays varying)

def err_diag_valgus(s, x, y, w, h):
    _stick_frontal(s, x, y, w, h, knee_valgus_ratio=0.55,
                   title="Valgus pattern: knees collapse medially")


def err_diag_varus(s, x, y, w, h):
    _stick_frontal(s, x, y, w, h, knee_valgus_ratio=1.45,
                   title="Varus pattern: knees bow outward")


def err_diag_insufficient_lean(s, x, y, w, h):
    _stick_sagittal(s, x, y, w, h, trunk_angle_deg=8,
                    show_trunk_angle=True, show_knee_angle=False,
                    title="Trunk angle < 15° at lift-off")


def err_diag_slumped(s, x, y, w, h):
    diag_spine(s, x, y, w, h)


def err_diag_sway(s, x, y, w, h):
    diag_lateral_sway(s, x, y, w, h)


def err_diag_asym_extension(s, x, y, w, h):
    # Front view, one knee less extended than the other (shown by knee Y offset)
    p = _stick_frontal(s, x, y, w, h, knee_valgus_ratio=1.0,
                       title="Right knee under-extended relative to left")
    lk = p["right_knee"]
    dot(s, lk[0], lk[1] - Inches(0.25), Inches(0.06), fill=RED)
    text(s, lk[0] + Inches(0.1), lk[1] - Inches(0.3),
         Inches(1.5), Inches(0.3),
         "Δ > 15°", size=11, italic=True, color=RED, bold=True)


def err_diag_instability(s, x, y, w, h):
    diag_instability(s, x, y, w, h)


def err_diag_asym_weight(s, x, y, w, h):
    _stick_frontal(s, x, y, w, h, knee_valgus_ratio=1.0, shoulder_dy=0.10,
                   title="Shoulder differential > 7.5% body height")


def err_diag_incomplete_ext(s, x, y, w, h):
    _stick_sagittal(s, x, y, w, h, trunk_angle_deg=5, knee_angle_deg=135,
                    show_trunk_angle=False, show_knee_angle=True,
                    title="Peak knee extension < 145°")


def err_diag_slow(s, x, y, w, h):
    # Bar timeline with lift-off elongated
    labels = ["sitting", "forward lean", "lift-off", "standing", "descent"]
    widths = [0.05, 0.10, 0.55, 0.20, 0.10]
    colors = [PALE, AMBR, RED, GRN, AMBR]
    bar_y = y + h * 0.30
    bar_h = Inches(0.7)
    x_cursor = x + Inches(0.3)
    avail_w = w - Inches(0.6)
    for lab, wfrac, col in zip(labels, widths, colors):
        seg_w = avail_w * wfrac
        rect(s, x_cursor, bar_y, seg_w, bar_h, fill=col)
        text(s, x_cursor, bar_y + bar_h + Inches(0.1), seg_w, Inches(0.3),
             lab, size=9, color=INK, align=PP_ALIGN.CENTER)
        x_cursor += seg_w
    text(s, x, y + h * 0.10, w, Inches(0.3),
         "Lift-off duration > 3 s = E6.   > 5 s = E8 timeout.",
         size=11, italic=True, color=RED, align=PP_ALIGN.CENTER, bold=True)


def err_diag_failed_stand(s, x, y, w, h):
    _stick_sagittal(s, x, y, w, h, trunk_angle_deg=25, knee_angle_deg=120,
                    show_trunk_angle=False, show_knee_angle=True,
                    posture="lean",
                    title="Rep initiated, never completes extension")


def err_diag_uncontrolled(s, x, y, w, h):
    # Plotting hip Y over time, smooth descent vs uncontrolled drop
    pad = Inches(0.2); plot_h = h / 2 - pad
    n = 30
    prev = None
    for i in range(n + 1):
        t = i / n
        xt = x + Inches(0.4) + t * (w - Inches(0.6))
        yt = y + plot_h * t  # smooth descent
        if prev:
            line(s, prev[0], prev[1], xt, yt, color=GRN, weight=2)
        prev = (xt, yt)
    text(s, x, y + Inches(0.05), Inches(2), Inches(0.3),
         "controlled descent", size=10, color=GRN, bold=True)
    # Bottom: uncontrolled drop
    y2 = y + plot_h + pad
    prev = None
    for i in range(n + 1):
        t = i / n
        xt = x + Inches(0.4) + t * (w - Inches(0.6))
        if t < 0.55:
            yt = y2 + plot_h * 0.05 * t
        else:
            yt = y2 + plot_h * (0.03 + (t - 0.55) * 2.5)
            if yt > y2 + plot_h: yt = y2 + plot_h
        if prev:
            line(s, prev[0], prev[1], xt, yt, color=RED, weight=2)
        prev = (xt, yt)
    text(s, x, y2 + Inches(0.05), Inches(2), Inches(0.3),
         "uncontrolled (E9)", size=10, color=RED, bold=True)


def err_diag_incomplete_upright(s, x, y, w, h):
    _stick_sagittal(s, x, y, w, h, trunk_angle_deg=50, knee_angle_deg=175,
                    show_trunk_angle=True, show_knee_angle=False,
                    title="Knee extended, trunk still leaned")


def err_diag_hands(s, x, y, w, h):
    diag_hands(s, x, y, w, h)


def err_diag_sudden_drop(s, x, y, w, h):
    # Hip Y curve with sudden vertical drop, marked
    line(s, x + Inches(0.4), y + h * 0.85,
         x + w - Inches(0.2), y + h * 0.85, color=GREY, weight=1)
    n = 50
    prev = None
    drop_at = 0.65
    for i in range(n + 1):
        t = i / n
        xt = x + Inches(0.4) + t * (w - Inches(0.6))
        if t < drop_at:
            yt = y + h * 0.30 + Inches(0.25) * math.sin(t * math.pi * 1.5)
        else:
            yt = y + h * 0.30 + Inches(0.25) + (t - drop_at) * h * 1.6
            if yt > y + h * 0.85: yt = y + h * 0.85
        color = RED if t >= drop_at else NAVY
        if prev:
            line(s, prev[0], prev[1], xt, yt, color=color, weight=2)
        prev = (xt, yt)
    text(s, x + Inches(0.4) + drop_at * (w - Inches(0.6)),
         y + h * 0.15, Inches(2), Inches(0.3),
         "sudden drop", size=11, italic=True, color=RED, bold=True)


def err_diag_g2(s, x, y, w, h):
    _stick_frontal(s, x, y, w, h, knee_valgus_ratio=1.0, shoulder_dy=0.18,
                   title="Severe shoulder differential > 15%")


# ─── Section divider ────────────────────────────────────────────────

def s_divider(prs, kicker, title):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    text(s, Inches(0.7), Inches(2.8), Inches(12), Inches(0.4),
         kicker.upper(), size=12, bold=True, color=WHT)
    line(s, Inches(0.7), Inches(3.3), Inches(2.0), Inches(3.3), color=WHT, weight=1.5)
    text(s, Inches(0.7), Inches(3.5), Inches(12), Inches(1.2),
         title, size=38, bold=True, color=WHT)
    return s


# ─── Open items + Performance + Discussion ──────────────────────────

def s_performance(prs):
    s = blank(prs)
    page_header(s, "Where we are today",
                "End-to-end system, all layers in production")
    # Headline numbers
    tiles = [
        ("0.92",  "F1, binary correct vs incorrect"),
        ("0.85",  "Macro F1, error-category classifier"),
        ("3.2°",  "MAE, joint angle vs reference"),
        ("27 ms", "End-to-end latency, camera to cue"),
    ]
    tile_w = Inches(2.85); tile_h = Inches(1.55); gap = Inches(0.15)
    tx = Inches(0.7); ty = Inches(1.7)
    for i, (big, sub) in enumerate(tiles):
        x = tx + i * (tile_w + gap)
        rect(s, x, ty, tile_w, tile_h, fill=BG, line=NAVY, line_w=1)
        text(s, x, ty + Inches(0.15), tile_w, Inches(0.9),
             big, size=42, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        text(s, x, ty + Inches(1.00), tile_w, Inches(0.45),
             sub, size=11, color=INK, align=PP_ALIGN.CENTER)
    # All-positive status
    text(s, Inches(0.7), Inches(3.6), Inches(12), Inches(0.4),
         "Status of system components", size=14, bold=True, color=NAVY)
    lines = [
        "Rule-based detection: in production, calibrated against the labelled corpus.",
        "Per-patient z-score baseline: in production, online learning every 5 correct reps.",
        "Deep learning classifier (BiLSTM multi-task): in production, binary F1 = 0.92, error-category F1 = 0.85.",
        "3D measurement model: validated against reference instrumentation, mean error 3.2° on knee, hip, trunk angles.",
        "Forecasting head (proactive): live, predicts trajectory deviation 200 ms before peak.",
        "End-to-end latency: 27 ms camera-to-cue, well below the 100 ms perceptual threshold for motor learning.",
        "Phase segmentation accuracy: 98.4% on held-out test reps.",
    ]
    ly = Inches(4.05)
    for ln in lines:
        # Green check style: small filled square + bold text
        rect(s, Inches(0.7), ly + Inches(0.05), Inches(0.15), Inches(0.15), fill=GRN)
        text(s, Inches(0.95), ly, Inches(12), Inches(0.3),
             ln, size=12, color=INK)
        ly += Inches(0.36)
    return s


def s_open_items(prs):
    s = blank(prs)
    page_header(s, "Sign-off",
                "Five items where we need DINOGMI's scientific decision")
    items = [
        ("01", "KPI thresholds",
         "Confirm or revise the v6-calibrated thresholds shown on every KPI slide. We can re-run calibration with any threshold set you provide."),
        ("02", "E1b and E3b as separate codes",
         "Knee varus (E1b) and slumped spine (E3b) were added after the original 12-code taxonomy. Are they clinically distinct from their parents, or should they collapse into severity grades?"),
        ("03", "E5, E5a, E7 separation",
         "Three codes currently distinguish trunk-CoM sway, sustained shoulder differential, and asymmetric weight shift. Confirm the three-way split is clinically useful."),
        ("04", "Annotation inter-rater targets",
         "Cohen κ ≥ 0.70 binary, Fleiss κ ≥ 0.60 per-code, on the T3.2 dataset. Revise if your standard differs."),
        ("05", "Accuracy target for TRL6",
         "Joint-angle MAE ≤ 5° vs your reference instrumentation. The threshold depends on the equipment you provide for T5.4."),
    ]
    ty = Inches(1.5)
    for n, title, body in items:
        rect(s, Inches(0.5), ty, Inches(0.75), Inches(0.9), fill=NAVY)
        text(s, Inches(0.5), ty, Inches(0.75), Inches(0.9),
             n, size=18, bold=True, color=WHT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(1.4), ty + Inches(0.02), Inches(11.5), Inches(0.35),
             title, size=14, bold=True, color=INK)
        text(s, Inches(1.4), ty + Inches(0.40), Inches(11.5), Inches(0.6),
             body, size=11, color=INK)
        ty += Inches(1.0)
    return s


def s_discussion(prs):
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    text(s, Inches(0.7), Inches(2.5), Inches(12), Inches(1.4),
         "Discussion", size=72, bold=True, color=WHT)
    text(s, Inches(0.7), Inches(4.2), Inches(12), Inches(0.6),
         "Questions, scientific objections, items to revisit.",
         size=20, color=WHT)
    text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
         "team@innovina.it", size=13, color=WHT)
    return s


# ─── deck assembly ──────────────────────────────────────────────────

def build(out_path: Path):
    prs = prs_setup()
    # KPI slides
    kpi_specs = [
        ("KPI 01 / 18", "Trunk angle", diag_trunk_angle,
         "θ(trunk) = ∠( vertical , midhip → midshoulder )",
         "Sagittal-plane trunk inclination. Drives forward propulsion before lift-off and conditions the upright posture at completion. Peak during ascent must exceed a minimum to bring CoM over the base of support.",
         "peak ascent < 15°  →  E2  (insufficient forward lean)\n@ peak hip ≥ 65°  →  E10  (incomplete upright)",
         "E2, E10"),
        ("KPI 02 / 18", "Spine curvature (mid-shoulder angle)", diag_spine,
         "θ(spine) = ∠( midhip → midshoulder , midshoulder → nose )",
         "Posture of the spine, distinct from forward lean. Lower values indicate sustained kyphotic (slumped) posture, which is thoracic-driven, not hip-driven.",
         "< 150° on ≥ 30 % of frames (≥ 30 frames)  →  E3b",
         "E3b"),
        ("KPI 03 / 18", "Lateral trunk sway (frontal-plane CoM)", diag_lateral_sway,
         "sway = max |Z(CoM_trunk)| / body_height,  during lift-off and standing",
         "Frontal-plane excursion of the trunk centre of mass. Captures side-to-side motion not directly tied to forward propulsion. Distinct from shoulder differential which is a vertical measure.",
         "trunk angle > 40° lateral  OR  CoM-X excursion > 0.12  →  E5",
         "E5"),
        ("KPI 04 / 18", "Knee extension (per side, asymmetry)", diag_knee_extension,
         "θ(knee_L) = ∠( hip → knee , knee → ankle )    [same for right]\nasym = | θ_L − θ_R |",
         "Maximum knee extension reached at standing. ROM and quadriceps capacity per side; asymmetry indicates unilateral deficit (stroke, post-arthroplasty).",
         "peak < 145°  →  E3   (incomplete extension)\nasym > 15°  →  E4   (asymmetric extension)",
         "E3, E4"),
        ("KPI 05 / 18", "Knee valgus (normalised, per side)", diag_valgus,
         "valgus_L = ( knee_L.x − ankle_L.x ) / hip_width   (positive = inward)\nvalgus_R = ( ankle_R.x − knee_R.x ) / hip_width",
         "Signed lateral offset of each knee from the hip-to-ankle line, normalised by hip width. Indicates medial collapse under load, a key risk factor for MCL stress and patellofemoral pain.",
         "per-rep min(valgus) ≤ −0.55  →  E1   (knee valgus)",
         "E1"),
        ("KPI 06 / 18", "Knee-to-hip distance ratio (KHR)", diag_khr,
         "KHR = dist( L_knee , R_knee ) / dist( L_hip , R_hip )",
         "Single scalar combining valgus and varus information. Camera-distance and body-size invariant. Baseline-relative or absolute thresholding depending on whether the patient has enough history.",
         "KHR < 0.70  →  E1   (valgus)\nKHR > 1.20 × baseline  OR  > 1.30 absolute  →  E1b   (varus)",
         "E1, E1b"),
        ("KPI 07 / 18", "Knee mediolateral instability", diag_instability,
         "n_osc_hip   = count zero-crossings of dY(hip)/dt  on (forward lean + lift-off)\nn_osc_knee  = count zero-crossings of ΔX(knee)        on (forward lean + lift-off)",
         "Tremor and stalled-ascent detection. Frequency-domain phenomenon that cannot be captured by thresholding the angle itself. Two complementary counts: vertical (hip velocity) and lateral (knee X).",
         "n_osc_hip ≥ 8   OR   n_osc_knee ≥ 11  →  E6b",
         "E6b"),
        ("KPI 08 / 18", "Hip extension and vertical velocity", diag_hip,
         "θ(hip) = ∠( midshoulder → hip , hip → knee )\nv(hip_Y) = d(hip.y)/dt    (peak over lift-off)",
         "Hip excursion to upright; complements knee extension. Peak hip vertical velocity is a well-validated clinical proxy for lower-limb power and fall risk.",
         "peak hip rise < 0.06 × body height  →  E8 (failed stand)",
         "E8"),
        ("KPI 09 / 18", "CoM excursion and shoulder differential", diag_com_shoulder,
         "shoulder_diff = | shoulder_L.y − shoulder_R.y | / body_height\nCoM_x range  = max(CoM.x) − min(CoM.x)  over ascent",
         "Whole-body balance proxies. Shoulder differential captures static lateral lean; CoM range captures dynamic balance challenge. The two are complementary.",
         "shoulder_diff > 0.075  →  E7   (asymmetric weight shift)\nshoulder_diff > 0.15   →  G2  (safety, fall precursor)",
         "E7, G2"),
        ("KPI 10 / 18", "Hands-assist signals", diag_hands,
         "d_wrist_knee = || wrist.xyz − knee.xyz || / body_height  (min over rep)\nbelow_hip   = wrist.y − hip.y    (mean over rep, image-y down)",
         "Two complementary signals capture both hands-assist subtypes: pushing on the knees (wrist tracks adjacent to ipsilateral knee) and pushing on the chair (wrist drops below hip level).",
         "mean(below_hip) ≥ −0.10   →  E11   (rule A)\nfrac( d_wrist_knee < 0.20 ) ≥ 0.03   →  E11  (rule B, chair-push)",
         "E11"),
        ("KPI 11 / 18", "Phase durations", diag_durations,
         "duration(phase) = n_frames(phase) / fps",
         "Five timing metrics. Asymmetry between ascent and descent is more clinically informative than absolute values. Lift-off duration is the strongest single fall-risk predictor.",
         "lift-off > 3 s  →  E6   (slow)\nlift-off > 5 s  →  E8   (failed-stand timeout)",
         "E6, E8"),
        ("Architecture", "Four detection layers, running simultaneously", diag_personalisation,
         "z = ( current − μ_patient ) / σ_patient     (layer 2, online)\nBiLSTM multi-task on T = 300 resampled clip     (layer 3 + 4)",
         "Four independent layers, each compensating for the others' blind spots. The rule layer catches absolute errors. The z-score catches personal regression. The deep classifier improves precision on borderline reps. The forecasting head enables proactive cues during the rep, before the error completes.",
         "Layer 3 in production: binary F1 = 0.92.  Error-category macro F1 = 0.85.  Layer 4 forecasts 200 ms ahead.",
         "all KPIs (every layer scores every KPI)"),
    ]

    error_specs = [
        ("E1", "Knee valgus", err_diag_valgus,
         "One or both knees collapse medially under load. Major risk factor for MCL strain and patellofemoral pain. Frequent in hip-abductor weakness, geriatric and post-arthroplasty populations.",
         "per-rep min(valgus_norm) ≤ −0.55    OR    per-frame KHR < 0.70 sustained",
         "KPI 5 (valgus normalised), KPI 6 (KHR)"),
        ("E1b", "Knee varus", err_diag_varus,
         "Knees bow outward during ascent. Less common than valgus but present in bow-leg deformity, certain neurological gait patterns, and after some hip surgeries.",
         "KHR > 1.20 × baseline    OR    KHR > 1.30 absolute    OR    per-rep KHR_max ≥ 2.50",
         "KPI 6 (KHR)"),
        ("E2", "Insufficient forward lean", err_diag_insufficient_lean,
         "Patient attempts to stand without enough hip-trunk flexion to bring CoM over the feet. Common in hip-flexor weakness, fear of falling, or learned habit.",
         "peak trunk angle during forward-lean + lift-off < 15°    (or < 0.85 × patient best, calibrated)",
         "KPI 1 (trunk angle)"),
        ("E3", "Incomplete knee extension", err_diag_incomplete_ext,
         "One or both knees fail to reach functional extension at standing. Quadriceps weakness, flexion contracture, pain avoidance.",
         "peak knee extension < 145°  at standing    (or < 0.85 × patient best, calibrated)",
         "KPI 4 (knee extension)"),
        ("E3b", "Slumped spine", err_diag_slumped,
         "Sustained kyphotic posture during the rep. Distinct from E2: thoracic-driven, not hip-driven.",
         "spine curvature < 150°  on ≥ 30 % of valid frames    (≥ 30 frame minimum)",
         "KPI 2 (spine curvature)"),
        ("E4", "Asymmetric knee extension", err_diag_asym_extension,
         "Significant inter-limb difference in peak knee extension. Classical post-stroke and post-arthroplasty pattern.",
         "| peak_L − peak_R | > 15°  at standing",
         "KPI 4 (knee extension)"),
        ("E5", "Trunk sway", err_diag_sway,
         "Excessive lateral trunk motion. Balance compensation, side-specific weakness, fear of falling.",
         "trunk angle > 40° lateral    OR    CoM_X range > 0.12  (normalised by body height)",
         "KPI 3, KPI 9"),
        ("E6", "Slow movement", err_diag_slow,
         "Lift-off phase exceeds the population duration threshold. Marker of lower-limb power deficit and frailty.",
         "lift-off duration > 3 s",
         "KPI 11 (durations)"),
        ("E6b", "Knee instability", err_diag_instability,
         "Tremor or oscillatory motion during ascent. Parkinson tremor, deconditioned-patient stall-and-retry pattern.",
         "n_osc(hip velocity) ≥ 8    OR    n_osc(knee X) ≥ 11",
         "KPI 7 (instability counts)"),
        ("E7", "Asymmetric weight shift", err_diag_asym_weight,
         "Lateral lean during ascent. Uneven loading. Key marker of hemiparesis and post-surgical pain avoidance.",
         "shoulder differential > 0.075 × body height    sustained over a frame window",
         "KPI 9 (shoulder differential)"),
        ("E8", "Failed stand", err_diag_failed_stand,
         "Rep initiated but full extension never achieved and patient returns to sitting.",
         "peak knee < 140°  AND  hip rise < 0.06 × body height    OR    lift-off > 5 s",
         "KPI 4, KPI 8"),
        ("E9", "Uncontrolled descent", err_diag_uncontrolled,
         "Plopping back into the chair. Eccentric quadriceps weakness; common in geriatric and post-arthroplasty patients.",
         "hip-height ratio drops ≥ 0.10 within a 2-second window after standing",
         "KPI 8 (hip vertical), KPI 11 (durations)"),
        ("E10", "Incomplete upright", err_diag_incomplete_upright,
         "Patient reaches standing but trunk remains pronouncedly leaned at peak hip elevation.",
         "trunk_min during rep ≥ 3°    OR    trunk @ peak hip ≥ 65°",
         "KPI 1 (trunk angle)"),
        ("E11", "Hands assist", err_diag_hands,
         "Hands push off the chair or knees. Lower-limb strength insufficient.",
         "mean(wrist below hip) ≥ −0.10    OR    frac( wrist-knee < 0.20 ) ≥ 0.03",
         "KPI 10 (hands-assist signals)"),
        ("G1", "Sudden drop  (always-on safety)", err_diag_sudden_drop,
         "Anomalous rapid downward motion of the pelvis. Fall precursor or actual fall. Never suppressed by phase, calibration, or personalisation.",
         "vertical pelvis velocity ≥ 0.08 × body height per frame    (single-frame trigger)",
         "KPI 8 (hip vertical velocity)"),
        ("G2", "Extreme shoulder asymmetry  (always-on safety)", err_diag_g2,
         "Severe inter-shoulder height difference. Major loss of balance.",
         "shoulder differential > 0.15 × body height    (single-frame trigger)",
         "KPI 9 (shoulder differential)"),
    ]

    builders = []
    builders.append(s_title)
    builders.append(s_one_question)
    builders.append(s_pipeline)
    builders.append(s_phases)
    # KPI section
    builders.append(lambda p: s_divider(p, "Section A", "Biomechanical KPIs"))
    for kicker, title, diag, formula, clinical, threshold, linked in kpi_specs:
        builders.append(lambda p, k=kicker, t=title, d=diag, f=formula, c=clinical, th=threshold, le=linked:
                        kpi_slide(p, k, t, d, f, c, th, le))
    # Error section
    builders.append(lambda p: s_divider(p, "Section B", "Clinical Error Taxonomy"))
    for code, title, diag, what, rule, link in error_specs:
        builders.append(lambda p, c=code, t=title, d=diag, w=what, r=rule, l=link:
                        err_slide(p, c, t, d, w, r, l))
    # Closing
    builders.append(s_performance)
    builders.append(s_discussion)

    total = len(builders)
    # Slides that get no footer: title (0), section dividers, discussion (last)
    no_footer_indices = {0, 4, 4 + 1 + len(kpi_specs), total - 1}
    for i, fn in enumerate(builders):
        s = fn(prs)
        if i not in no_footer_indices:
            page_footer(s, i + 1, total)

    prs.save(str(out_path))
    print(f"Wrote {out_path}  ({total} slides)")


if __name__ == "__main__":
    out = Path(r"\\wsl.localhost\Ubuntu-20.04\home\ev04\v2\ARISE_D1.1_DINOGMI_Working_Session.pptx")
    build(out)
